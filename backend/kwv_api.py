# -*- coding: utf-8 -*-
"""
KoreaWorkingVisa API 모듈
- 관리자/비자신청자 인증 및 관리
- Google OAuth 연동
- JWT 토큰 기반 인증
"""

from fastapi import APIRouter, HTTPException, Depends, Header, UploadFile, File, Form, Request, Body
from fastapi.responses import JSONResponse, RedirectResponse, FileResponse
from typing import Optional, List
from pydantic import BaseModel
from datetime import datetime, timedelta, date
import os
import json
import hashlib
import secrets
import pymysql
import pymysql.cursors
import uuid
import base64
import httpx

# JWT 관련 (python-jose)
try:
    from jose import JWTError, jwt
    JWT_AVAILABLE = True
except ImportError:
    JWT_AVAILABLE = False
    print("⚠️ python-jose not installed. JWT features will be limited.")

# bcrypt 관련
try:
    import bcrypt
    BCRYPT_AVAILABLE = True
except ImportError:
    BCRYPT_AVAILABLE = False
    print("⚠️ bcrypt not installed. Using simple hash fallback.")

# ==================== 설정 ====================
JWT_SECRET_KEY = os.getenv("JWT_SECRET_KEY", "kwv-secret-key-change-in-production-2026")
JWT_ALGORITHM = os.getenv("JWT_ALGORITHM", "HS256")
JWT_EXPIRATION_HOURS = int(os.getenv("JWT_EXPIRATION_HOURS", "24"))

GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET", "")
SITE_URL = os.getenv("SITE_URL", "")  # 예: http://114.202.247.46:8080

# 파일 업로드 경로
LOCAL_UPLOAD_DIR = os.path.join(os.path.dirname(__file__), "file_uploads")
os.makedirs(LOCAL_UPLOAD_DIR, exist_ok=True)

# 기본 비밀번호
DEFAULT_PASSWORD = "kwv2026"

# ==================== Router ====================
router = APIRouter(prefix="/api/kwv", tags=["KoreaWorkingVisa"])

APP_VERSION = "0.2.20260305"

@router.get("/version")
async def get_version():
    """앱 버전 정보"""
    return {"version": APP_VERSION}

@router.get("/readme")
async def get_readme():
    """README.md 내용 반환"""
    readme_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "README.md")
    if os.path.exists(readme_path):
        with open(readme_path, 'r', encoding='utf-8') as f:
            return {"content": f.read()}
    return {"content": "README.md not found"}

# ==================== Pydantic Models ====================
class UserLogin(BaseModel):
    email: str
    password: str

class UserRegister(BaseModel):
    email: str
    name: str
    phone: Optional[str] = None
    address: Optional[str] = None
    language: Optional[str] = "en"
    user_type: Optional[str] = "applicant"  # applicant 또는 admin
    organization: Optional[str] = None  # 관리자: 소속 지역
    profile_photo: Optional[str] = None  # Base64 또는 URL
    passport_copy_url: Optional[str] = None
    visa_copy_url: Optional[str] = None
    id_card_url: Optional[str] = None
    insurance_cert_url: Optional[str] = None
    # Phase 2 추가
    nationality: Optional[str] = None  # 국적
    visa_type: Optional[str] = None  # 비자 유형 (E-8, E-9 등)
    birth_date: Optional[str] = None  # 생년월일
    gender: Optional[str] = None  # 성별
    target_local_government_id: Optional[int] = None  # 신청 대상 지자체

class AdminLoginRequest(BaseModel):
    name: str
    password: str

class AdminRegisterRequest(BaseModel):
    name: str
    birth_date: str  # YYYY-MM-DD
    phone: str
    organization: str

class GoogleLoginRequest(BaseModel):
    credential: Optional[str] = None  # Google ID token (legacy)
    email: Optional[str] = None
    name: Optional[str] = None

class TokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    expires_in: int
    user: dict

class ApplicantCreate(BaseModel):
    nationality: str
    passport_number: str
    birth_date: str
    gender: str
    visa_type: str
    employer_name: Optional[str] = None
    job_category: Optional[str] = None

class ApplicantStatusUpdate(BaseModel):
    status: str  # pending, processing, approved, rejected
    rejection_reason: Optional[str] = None

# ==================== 유틸리티 함수 ====================

def hash_password(password: str) -> str:
    """비밀번호 해시"""
    if BCRYPT_AVAILABLE:
        return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
    else:
        return hashlib.sha256(password.encode()).hexdigest()

def verify_password(password: str, hashed: str) -> bool:
    """비밀번호 검증"""
    if BCRYPT_AVAILABLE:
        try:
            return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))
        except:
            return False
    else:
        return hashlib.sha256(password.encode()).hexdigest() == hashed

# DB 세션 토큰 저장소 (JWT 불가 시 사용)
_token_store = {}

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None) -> str:
    """JWT 토큰 생성 (JWT 불가 시 DB 세션 토큰)"""
    if not JWT_AVAILABLE:
        import secrets
        token = secrets.token_hex(32)
        expire = datetime.utcnow() + timedelta(hours=JWT_EXPIRATION_HOURS)
        _token_store[token] = {**data, "exp": expire.isoformat()}
        return token

    to_encode = data.copy()
    expire = datetime.utcnow() + (expires_delta or timedelta(hours=JWT_EXPIRATION_HOURS))
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, JWT_SECRET_KEY, algorithm=JWT_ALGORITHM)

def decode_token(token: str) -> Optional[dict]:
    """JWT 토큰 디코딩"""
    if not JWT_AVAILABLE:
        stored = _token_store.get(token)
        if not stored:
            return None
        if datetime.fromisoformat(stored["exp"]) < datetime.utcnow():
            del _token_store[token]
            return None
        return {k: v for k, v in stored.items() if k != "exp"}
    try:
        payload = jwt.decode(token, JWT_SECRET_KEY, algorithms=[JWT_ALGORITHM])
        return payload
    except JWTError:
        return None

async def get_current_user(authorization: Optional[str] = Header(None)) -> dict:
    """현재 로그인한 사용자 정보 가져오기"""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="인증이 필요합니다")

    token = authorization.replace("Bearer ", "")
    payload = decode_token(token)

    if not payload:
        raise HTTPException(status_code=401, detail="유효하지 않은 토큰입니다")

    return payload

def require_admin(user: dict) -> dict:
    """관리자 권한 확인"""
    if user.get("user_type") != "admin":
        raise HTTPException(status_code=403, detail="관리자 권한이 필요합니다")
    return user

def require_admin_level(user: dict, min_level: int) -> dict:
    """특정 등급 이상의 관리자 권한 확인 (숫자가 클수록 높은 권한, 9=super admin)"""
    if user.get("user_type") != "admin":
        raise HTTPException(status_code=403, detail="관리자 권한이 필요합니다")
    if user.get("admin_level", 0) < min_level:
        raise HTTPException(status_code=403, detail=f"등급 {min_level} 이상의 관리자 권한이 필요합니다")
    return user

# ==================== 데이터베이스 연결 ====================
MOCK_MODE = os.getenv('KWV_MOCK_MODE', 'false').lower() == 'true'

# Mock 데이터 (MOCK_MODE=true일 때만 사용)
MOCK_USERS = {}
MOCK_USER_ID_COUNTER = 1
MOCK_APPLICATIONS = []

def get_kwv_db_connection():
    """KWV 데이터베이스 연결"""
    if MOCK_MODE:
        return None

    import pymysql
    from pathlib import Path
    from dotenv import load_dotenv

    env_path = Path(__file__).parent.parent / '.env'
    load_dotenv(dotenv_path=env_path, override=True)

    try:
        return pymysql.connect(
            host=os.getenv('DB_HOST', 'localhost'),
            user=os.getenv('DB_USER', 'root'),
            passwd=os.getenv('DB_PASSWORD', ''),
            db=os.getenv('DB_NAME', 'koreaworkingvisa'),
            charset='utf8mb4',
            port=int(os.getenv('DB_PORT', '3306'))
        )
    except Exception as e:
        print(f"⚠️ Database connection failed: {e}")
        return None

# ==================== 인증 API ====================

@router.post("/auth/register")
async def register(user_data: UserRegister):
    """
    회원가입
    - applicant (일반 사용자): 비자 신청자 (여권+비자 첨부)
    - admin (관리자): 지역별 담당자 (신분증+4대보험 첨부)
    - 비밀번호는 입력받지 않고 기본값 'kwv2026' 사용
    """
    global MOCK_USER_ID_COUNTER

    if user_data.user_type not in ["applicant", "admin"]:
        raise HTTPException(status_code=400, detail="유효하지 않은 사용자 유형입니다")

    # 프로필 사진 Base64 → 파일 저장
    profile_photo_url = None
    if user_data.profile_photo and user_data.profile_photo.startswith('data:'):
        profile_photo_url = save_base64_file(user_data.profile_photo, "profile")
    elif user_data.profile_photo:
        profile_photo_url = user_data.profile_photo

    # Mock mode 처리
    if MOCK_MODE:
        for u in MOCK_USERS.values():
            if u["email"] == user_data.email:
                raise HTTPException(status_code=400, detail="이미 등록된 이메일입니다")

        user_id = MOCK_USER_ID_COUNTER
        MOCK_USER_ID_COUNTER += 1

        is_admin = user_data.user_type == "admin"
        new_user = {
            "id": user_id,
            "email": user_data.email,
            "password_hash": hash_password(DEFAULT_PASSWORD),
            "name": user_data.name,
            "phone": user_data.phone,
            "address": user_data.address,
            "user_type": user_data.user_type,
            "admin_level": 2 if is_admin else None,
            "language": user_data.language,
            "is_active": True,
            "is_approved": True,
            "organization": user_data.organization,
            "profile_photo": profile_photo_url,
            "created_at": datetime.utcnow().isoformat()
        }
        MOCK_USERS[user_id] = new_user

        token_data = {
            "sub": str(user_id),
            "email": user_data.email,
            "name": user_data.name,
            "user_type": user_data.user_type,
            "admin_level": 2 if is_admin else None
        }
        access_token = create_access_token(token_data)

        return {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user_id,
                "email": user_data.email,
                "name": user_data.name,
                "user_type": user_data.user_type,
                "admin_level": 2 if is_admin else None,
                "language": user_data.language
            }
        }

    # 실제 DB 모드
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()

        cursor.execute("SELECT id FROM kwv_users WHERE email = %s", (user_data.email,))
        if cursor.fetchone():
            raise HTTPException(status_code=400, detail="이미 등록된 이메일입니다")

        # 기본 비밀번호 kwv2026 해시
        password_salt = secrets.token_hex(16)
        password_hash = hashlib.sha256((DEFAULT_PASSWORD + password_salt).encode()).hexdigest()

        user_type = user_data.user_type or 'applicant'

        # 승인 모드 확인
        approval_mode = 'manual'
        try:
            cursor.execute("SELECT setting_value FROM kwv_system_settings WHERE setting_key = 'approval_mode'")
            row = cursor.fetchone()
            if row:
                approval_mode = row[0]
        except:
            pass

        # 모든 사용자는 미승인 상태로 생성 (관리자도 승인 필수)
        is_approved = False

        cursor.execute("""
            INSERT INTO kwv_users (email, password_hash, password_salt, name, phone, address,
                user_type, admin_level, language, organization, region, profile_photo,
                target_local_government_id, is_approved)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            user_data.email,
            password_hash,
            password_salt,
            user_data.name,
            user_data.phone or '',
            user_data.address or '',
            user_type,
            2 if user_type == 'admin' else 0,
            user_data.language or 'en',
            user_data.organization if user_type == 'admin' else None,
            user_data.organization if user_type == 'admin' else None,
            profile_photo_url,
            user_data.target_local_government_id if user_type == 'applicant' else None,
            is_approved
        ))

        user_id = cursor.lastrowid

        # 첨부 파일 정보 저장
        file_entries = []
        if user_data.passport_copy_url:
            file_entries.append((user_id, 'passport_copy', 'passport', user_data.passport_copy_url))
        if user_data.visa_copy_url:
            file_entries.append((user_id, 'visa_copy', 'visa', user_data.visa_copy_url))
        if user_data.id_card_url:
            file_entries.append((user_id, 'id_card', 'id_card', user_data.id_card_url))
        if user_data.insurance_cert_url:
            file_entries.append((user_id, 'insurance_cert', 'insurance', user_data.insurance_cert_url))

        for uid, category, fname, fpath in file_entries:
            cursor.execute("""
                INSERT INTO kwv_file_uploads (user_id, file_category, file_name, file_path)
                VALUES (%s, %s, %s, %s)
            """, (uid, category, fname, fpath))

        # applicant인 경우 비자 신청 정보 저장
        if user_type == 'applicant':
            cursor.execute("""
                INSERT INTO kwv_visa_applicants (user_id, visa_type, nationality, birth_date, gender)
                VALUES (%s, %s, %s, %s, %s)
            """, (
                user_id,
                user_data.visa_type,
                user_data.nationality,
                user_data.birth_date or None,
                user_data.gender or None
            ))

        # 자동 승인 모드일 때 필수 항목 검증
        auto_approved = False
        if approval_mode == 'auto' and user_type == 'applicant':
            check = check_auto_approval(user_id, conn)
            if check["passed"]:
                cursor.execute("""
                    UPDATE kwv_users SET is_approved = TRUE, approved_at = NOW() WHERE id = %s
                """, (user_id,))
                cursor.execute("""
                    UPDATE kwv_visa_applicants SET application_status = 'approved' WHERE user_id = %s
                """, (user_id,))
                auto_approved = True

        conn.commit()

        token_data = {
            "sub": str(user_id),
            "email": user_data.email,
            "name": user_data.name,
            "user_type": user_type,
            "admin_level": 2 if user_type == 'admin' else 0
        }
        access_token = create_access_token(token_data)

        response = {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user_id,
                "email": user_data.email,
                "name": user_data.name,
                "user_type": user_type,
                "admin_level": 2 if user_type == 'admin' else 0,
                "language": user_data.language,
                "is_approved": auto_approved if user_type == 'applicant' else False
            }
        }
        # 자동 승인 실패 시 누락 항목 알려주기
        if approval_mode == 'auto' and user_type == 'applicant' and not auto_approved:
            check = check_auto_approval(user_id, conn)
            response["approval_status"] = "pending"
            response["missing_items"] = check.get("missing", [])
        elif auto_approved:
            response["approval_status"] = "approved"

        return response
    except HTTPException:
        raise
    except Exception as e:
        if conn:
            conn.rollback()
        raise HTTPException(status_code=500, detail=f"회원가입 실패: {str(e)}")
    finally:
        if conn:
            conn.close()

@router.post("/auth/login", response_model=TokenResponse)
async def login(credentials: UserLogin):
    """일반 로그인 (이메일 + 비밀번호)"""

    if MOCK_MODE:
        user = None
        for u in MOCK_USERS.values():
            if u["email"] == credentials.email:
                user = u
                break

        if not user:
            raise HTTPException(status_code=401, detail="이메일 또는 비밀번호가 올바르지 않습니다")

        if not user.get("is_active", True):
            raise HTTPException(status_code=403, detail="비활성화된 계정입니다")

        if not verify_password(credentials.password, user["password_hash"]):
            raise HTTPException(status_code=401, detail="이메일 또는 비밀번호가 올바르지 않습니다")

        token_data = {
            "sub": str(user["id"]),
            "email": user["email"],
            "name": user["name"],
            "user_type": user["user_type"],
            "admin_level": user["admin_level"]
        }
        access_token = create_access_token(token_data)

        return {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user["id"],
                "email": user["email"],
                "name": user["name"],
                "user_type": user["user_type"],
                "admin_level": user["admin_level"],
                "language": user.get("language", "en")
            }
        }

    # 실제 DB 모드
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()
        # 이메일 또는 이름으로 로그인 가능
        login_input = credentials.email.strip()
        cursor.execute("""
            SELECT id, email, password_hash, password_salt, name, user_type, admin_level, language, is_active
            FROM kwv_users WHERE email = %s OR name = %s
        """, (login_input, login_input))

        user = cursor.fetchone()

        if not user:
            raise HTTPException(status_code=401, detail="이메일/이름 또는 비밀번호가 올바르지 않습니다")

        user_id, email, password_hash, password_salt, name, user_type, admin_level, language, is_active = user
        name = name or email
        status = 'active' if is_active else 'inactive'
        user_type = user_type or 'applicant'
        admin_level = admin_level or 0
        language = language or 'ko'

        if status != 'active':
            raise HTTPException(status_code=403, detail="비활성화된 계정입니다")

        check_hash = hashlib.sha256((credentials.password + password_salt).encode()).hexdigest()
        if check_hash != password_hash:
            raise HTTPException(status_code=401, detail="이메일 또는 비밀번호가 올바르지 않습니다")

        cursor.execute("UPDATE kwv_users SET last_login_at = NOW() WHERE id = %s", (user_id,))
        conn.commit()

        token_data = {
            "sub": str(user_id),
            "email": email,
            "name": name,
            "user_type": user_type,
            "admin_level": admin_level
        }
        access_token = create_access_token(token_data)

        return {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user_id,
                "email": email,
                "name": name,
                "user_type": user_type,
                "admin_level": admin_level,
                "language": language
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"로그인 실패: {str(e)}")
    finally:
        if conn:
            conn.close()

@router.post("/auth/google")
async def google_login(request: GoogleLoginRequest):
    """Google OAuth 로그인/가입 (credential 또는 email+name)"""

    email = None
    name = None
    google_id = None

    # 방법1: credential (Google ID token) - 레거시
    if request.credential:
        if not GOOGLE_CLIENT_ID:
            raise HTTPException(status_code=503, detail="Google OAuth가 설정되지 않았습니다")
        try:
            from google.oauth2 import id_token
            from google.auth.transport import requests as google_requests
            idinfo = id_token.verify_oauth2_token(
                request.credential, google_requests.Request(), GOOGLE_CLIENT_ID
            )
            email = idinfo.get('email')
            name = idinfo.get('name', email.split('@')[0])
            google_id = idinfo.get('sub')
        except ImportError:
            raise HTTPException(status_code=503, detail="Google Auth 라이브러리가 설치되지 않았습니다")
        except Exception as e:
            raise HTTPException(status_code=401, detail=f"Google 인증 실패: {str(e)}")
    # 방법2: email+name (콜백 경유)
    elif request.email:
        email = request.email
        name = request.name or email.split('@')[0]
    else:
        raise HTTPException(status_code=400, detail="credential 또는 email이 필요합니다")

    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()

        cursor.execute("SELECT id, name, user_type, admin_level, language FROM kwv_users WHERE email = %s", (email,))
        existing = cursor.fetchone()

        if existing:
            user_id, name, user_type, admin_level, language = existing
            if google_id:
                cursor.execute("UPDATE kwv_users SET last_login_at = NOW(), oauth_provider = 'google', oauth_id = %s WHERE id = %s",
                             (google_id, user_id))
            else:
                cursor.execute("UPDATE kwv_users SET last_login_at = NOW(), oauth_provider = 'google' WHERE id = %s", (user_id,))
        else:
            # 미가입 → 자동 가입 (applicant)
            cursor.execute("""
                INSERT INTO kwv_users (email, name, user_type, oauth_provider, oauth_id, language, is_approved)
                VALUES (%s, %s, 'applicant', 'google', %s, 'en', TRUE)
            """, (email, name, google_id))
            user_id = cursor.lastrowid
            user_type = 'applicant'
            admin_level = None
            language = 'en'

        conn.commit()

        token_data = {
            "sub": str(user_id),
            "email": email,
            "name": name,
            "user_type": user_type,
            "admin_level": admin_level
        }
        access_token = create_access_token(token_data)

        return {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user_id,
                "email": email,
                "name": name,
                "user_type": user_type,
                "admin_level": admin_level,
                "language": language
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Google 인증 처리 실패: {str(e)}")
    finally:
        conn.close()

@router.get("/auth/me")
async def get_me(user: dict = Depends(get_current_user)):
    """현재 로그인한 사용자 정보"""
    return {
        "id": user.get("sub"),
        "email": user.get("email"),
        "name": user.get("name"),
        "user_type": user.get("user_type"),
        "admin_level": user.get("admin_level")
    }

@router.post("/auth/logout")
async def logout():
    """로그아웃 (클라이언트에서 토큰 삭제)"""
    return {"message": "로그아웃 성공"}

# ==================== 관리자 API ====================

@router.get("/admin/applicants")
async def get_applicants(
    status: Optional[str] = None,
    nationality: Optional[str] = None,
    visa_type: Optional[str] = None,
    lg_id: Optional[int] = None,
    is_approved: Optional[str] = None,
    search: Optional[str] = None,
    page: int = 1,
    limit: int = 20,
    user: dict = Depends(get_current_user)
):
    """신청자 목록 조회 (관리자용) - 필터 강화"""
    require_admin(user)

    conn = get_kwv_db_connection()
    if not conn:
        return {"applicants": [], "total": 0, "page": page, "limit": limit}

    try:
        cursor = conn.cursor()

        where_clause = "WHERE u.user_type = 'applicant'"
        params = []

        if status:
            where_clause += " AND a.application_status = %s"
            params.append(status)
        if nationality:
            where_clause += " AND a.nationality = %s"
            params.append(nationality)
        if visa_type:
            where_clause += " AND a.visa_type = %s"
            params.append(visa_type)
        if lg_id:
            where_clause += " AND u.target_local_government_id = %s"
            params.append(lg_id)
        if is_approved == 'true':
            where_clause += " AND u.is_approved = TRUE"
        elif is_approved == 'false':
            where_clause += " AND u.is_approved = FALSE"
        if search:
            where_clause += " AND (u.name LIKE %s OR u.email LIKE %s OR u.phone LIKE %s)"
            params.extend([f"%{search}%", f"%{search}%", f"%{search}%"])

        cursor.execute(f"""
            SELECT COUNT(*) FROM kwv_users u
            LEFT JOIN kwv_visa_applicants a ON u.id = a.user_id
            {where_clause}
        """, params)
        total = cursor.fetchone()[0]

        offset = (page - 1) * limit
        cursor.execute(f"""
            SELECT u.id, u.email, u.name, u.phone, u.created_at,
                   a.visa_type, a.nationality, a.application_status,
                   u.is_approved, u.approved_at, u.target_local_government_id,
                   u.local_government_id, u.profile_photo, u.language,
                   lg.name as lg_name, tlg.name as target_lg_name,
                   a.birth_date, a.gender
            FROM kwv_users u
            LEFT JOIN kwv_visa_applicants a ON u.id = a.user_id
            LEFT JOIN kwv_local_governments lg ON u.local_government_id = lg.id
            LEFT JOIN kwv_local_governments tlg ON u.target_local_government_id = tlg.id
            {where_clause}
            ORDER BY u.created_at DESC
            LIMIT %s OFFSET %s
        """, params + [limit, offset])

        applicants = []
        for row in cursor.fetchall():
            applicants.append({
                "id": row[0],
                "email": row[1],
                "name": row[2],
                "phone": row[3],
                "created_at": row[4].isoformat() if row[4] else None,
                "visa_type": row[5],
                "nationality": row[6],
                "status": row[7] or "pending",
                "is_approved": bool(row[8]),
                "approved_at": row[9].isoformat() if row[9] else None,
                "target_local_government_id": row[10],
                "local_government_id": row[11],
                "profile_photo": row[12],
                "language": row[13],
                "lg_name": row[14],
                "target_lg_name": row[15],
                "birth_date": row[16].isoformat() if row[16] else None,
                "gender": row[17]
            })

        # 필터용 메타데이터
        cursor.execute("SELECT DISTINCT nationality FROM kwv_visa_applicants WHERE nationality IS NOT NULL ORDER BY nationality")
        nationalities = [r[0] for r in cursor.fetchall()]
        cursor.execute("SELECT DISTINCT visa_type FROM kwv_visa_applicants WHERE visa_type IS NOT NULL ORDER BY visa_type")
        visa_types = [r[0] for r in cursor.fetchall()]

        return {
            "applicants": applicants,
            "total": total,
            "page": page,
            "limit": limit,
            "filters": {
                "nationalities": nationalities,
                "visa_types": visa_types
            }
        }
    finally:
        conn.close()

@router.get("/admin/statistics")
async def get_statistics(user: dict = Depends(get_current_user)):
    """대시보드 통계 (관리자용)"""
    require_admin(user)

    conn = get_kwv_db_connection()
    if not conn:
        return {
            "total_applicants": 0,
            "pending": 0,
            "approved": 0,
            "rejected": 0,
            "processing": 0
        }

    try:
        cursor = conn.cursor()

        cursor.execute("SELECT COUNT(*) FROM kwv_users WHERE user_type = 'applicant'")
        total = cursor.fetchone()[0]

        stats = {"total_applicants": total, "pending": 0, "approved": 0, "rejected": 0, "processing": 0}

        cursor.execute("""
            SELECT application_status, COUNT(*) FROM kwv_visa_applicants
            GROUP BY application_status
        """)
        for row in cursor.fetchall():
            if row[0] in stats:
                stats[row[0]] = row[1]

        return stats
    finally:
        conn.close()

@router.put("/admin/applicants/{applicant_id}/status")
async def update_applicant_status(
    applicant_id: int,
    status_update: ApplicantStatusUpdate,
    user: dict = Depends(get_current_user)
):
    """신청자 상태 변경 (관리자용)"""
    require_admin_level(user, 2)

    valid_statuses = ["pending", "processing", "approved", "rejected"]
    if status_update.status not in valid_statuses:
        raise HTTPException(status_code=400, detail="유효하지 않은 상태입니다")

    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()

        cursor.execute("""
            UPDATE kwv_visa_applicants
            SET application_status = %s, rejection_reason = %s, updated_at = NOW()
            WHERE user_id = %s
        """, (status_update.status, status_update.rejection_reason, applicant_id))

        # kwv_users의 승인 상태도 동기화
        if status_update.status == 'approved':
            cursor.execute("""
                UPDATE kwv_users SET is_approved = TRUE, approved_at = NOW(),
                approved_by = %s, rejection_reason = NULL WHERE id = %s
            """, (user.get('sub'), applicant_id))
        elif status_update.status == 'rejected':
            cursor.execute("""
                UPDATE kwv_users SET is_approved = FALSE,
                rejection_reason = %s WHERE id = %s
            """, (status_update.rejection_reason, applicant_id))

        conn.commit()
        return {"message": "상태가 변경되었습니다", "status": status_update.status}
    finally:
        conn.close()

@router.put("/admin/applicants/{applicant_id}/assign-lg")
async def assign_local_government(applicant_id: int, request: Request, user: dict = Depends(get_current_user)):
    """근로자를 지자체에 배정"""
    require_admin_level(user, 2)
    body = await request.json()
    lg_id = body.get("local_government_id")
    if not lg_id:
        raise HTTPException(status_code=400, detail="지자체 ID가 필요합니다")

    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        # 지자체 TO 확인
        cursor.execute("SELECT allocated_quota, used_quota, name FROM kwv_local_governments WHERE id = %s AND is_active = TRUE", (lg_id,))
        lg = cursor.fetchone()
        if not lg:
            raise HTTPException(status_code=404, detail="지자체를 찾을 수 없습니다")
        if lg[0] > 0 and lg[1] >= lg[0]:
            raise HTTPException(status_code=400, detail=f"'{lg[2]}' 지자체의 TO가 부족합니다 ({lg[1]}/{lg[0]})")

        # 배정
        cursor.execute("UPDATE kwv_users SET local_government_id = %s WHERE id = %s AND user_type = 'applicant'", (lg_id, applicant_id))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="신청자를 찾을 수 없습니다")

        # used_quota 증가
        cursor.execute("UPDATE kwv_local_governments SET used_quota = used_quota + 1 WHERE id = %s", (lg_id,))
        conn.commit()
        return {"message": f"'{lg[2]}' 지자체에 배정되었습니다"}
    finally:
        conn.close()

def check_auto_approval(user_id: int, conn) -> dict:
    """자동 승인 필수 항목 검증 - 모든 항목이 충족되어야 자동 승인"""
    cursor = conn.cursor()
    missing = []

    # 1. 기본 정보 확인
    cursor.execute("""
        SELECT name, phone, profile_photo, target_local_government_id
        FROM kwv_users WHERE id = %s
    """, (user_id,))
    user_row = cursor.fetchone()
    if not user_row:
        return {"passed": False, "missing": ["사용자 정보 없음"]}

    name, phone, profile_photo, target_lg = user_row
    if not name: missing.append("이름")
    if not phone: missing.append("전화번호")
    if not profile_photo: missing.append("프로필 사진")
    if not target_lg: missing.append("신청 대상 지자체")

    # 2. 비자 정보 확인
    cursor.execute("""
        SELECT nationality, visa_type, passport_number, birth_date, gender
        FROM kwv_visa_applicants WHERE user_id = %s
    """, (user_id,))
    visa_row = cursor.fetchone()
    if not visa_row:
        missing.extend(["국적", "비자유형", "여권번호", "생년월일", "성별"])
    else:
        nat, vtype, passport, bdate, gender = visa_row
        if not nat: missing.append("국적")
        if not vtype: missing.append("비자유형")
        if not passport: missing.append("여권번호")
        if not bdate: missing.append("생년월일")
        if not gender: missing.append("성별")

    # 3. 필수 서류 확인 (여권 사본, 비자 사본)
    cursor.execute("""
        SELECT file_category FROM kwv_file_uploads WHERE user_id = %s
        AND file_category IN ('passport_copy', 'visa_copy')
    """, (user_id,))
    uploaded = {r[0] for r in cursor.fetchall()}
    if 'passport_copy' not in uploaded: missing.append("여권 사본")
    if 'visa_copy' not in uploaded: missing.append("비자 사본")

    return {"passed": len(missing) == 0, "missing": missing}

# ==================== 신청자 API ====================

@router.get("/my/profile")
async def get_my_profile(user: dict = Depends(get_current_user)):
    """내 프로필 조회"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"id": user.get("sub"), "email": user.get("email"), "name": user.get("name")}

    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, email, name, phone, language, profile_photo, created_at
            FROM kwv_users WHERE id = %s
        """, (user.get("sub"),))

        row = cursor.fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="사용자를 찾을 수 없습니다")

        return {
            "id": row[0],
            "email": row[1],
            "name": row[2],
            "phone": row[3],
            "language": row[4],
            "profile_photo": row[5],
            "created_at": row[6].isoformat() if row[6] else None
        }
    finally:
        conn.close()

@router.get("/my/application")
async def get_my_application(user: dict = Depends(get_current_user)):
    """내 비자 신청 현황"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"application": None}

    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT a.id, a.visa_type, a.nationality, a.passport_number,
                   a.application_status, a.created_at, a.updated_at
            FROM kwv_visa_applicants a
            WHERE a.user_id = %s
            ORDER BY a.created_at DESC LIMIT 1
        """, (user.get("sub"),))

        row = cursor.fetchone()
        if not row:
            return {"application": None}

        return {
            "application": {
                "id": row[0],
                "visa_type": row[1],
                "nationality": row[2],
                "passport_number": row[3],
                "status": row[4],
                "created_at": row[5].isoformat() if row[5] else None,
                "updated_at": row[6].isoformat() if row[6] else None
            }
        }
    finally:
        conn.close()

@router.post("/my/application")
async def create_application(
    application: ApplicantCreate,
    user: dict = Depends(get_current_user)
):
    """비자 신청"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()

        cursor.execute("SELECT id FROM kwv_visa_applicants WHERE user_id = %s", (user.get("sub"),))
        if cursor.fetchone():
            raise HTTPException(status_code=400, detail="이미 신청서가 존재합니다")

        cursor.execute("""
            INSERT INTO kwv_visa_applicants
            (user_id, visa_type, nationality, passport_number, birth_date, gender, employer_name, job_category)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            user.get("sub"),
            application.visa_type,
            application.nationality,
            application.passport_number,
            application.birth_date,
            application.gender,
            application.employer_name,
            application.job_category
        ))

        conn.commit()
        return {"message": "신청이 완료되었습니다", "application_id": cursor.lastrowid}
    finally:
        conn.close()

# ==================== 파일 업로드 API ====================

def upload_to_local(file_data: bytes, filename: str, category: str) -> str:
    """파일을 로컬에 저장하고 URL 반환"""
    cat_dir = os.path.join(LOCAL_UPLOAD_DIR, category)
    os.makedirs(cat_dir, exist_ok=True)
    file_path = os.path.join(cat_dir, filename)
    with open(file_path, 'wb') as f:
        f.write(file_data)
    return f"/api/kwv/uploads/{category}/{filename}"

def save_base64_file(base64_data: str, category: str) -> str:
    """Base64 데이터를 파일로 저장"""
    if ',' in base64_data:
        base64_data = base64_data.split(',', 1)[1]
    file_data = base64.b64decode(base64_data)
    ext = 'jpg'
    filename = f"{int(datetime.utcnow().timestamp())}_{uuid.uuid4().hex[:8]}.{ext}"
    return upload_to_local(file_data, filename, category)

@router.get("/uploads/{category}/{filename}")
async def serve_upload(category: str, filename: str):
    """업로드된 파일 서빙"""
    file_path = os.path.join(LOCAL_UPLOAD_DIR, category, filename)
    if os.path.exists(file_path):
        return FileResponse(file_path)
    raise HTTPException(status_code=404, detail="File not found")

@router.post("/auth/upload-temp")
async def upload_temp_file(file: UploadFile = File(...), category: str = Form("temp")):
    """가입 전 임시 파일 업로드"""
    allowed = {'image/jpeg', 'image/png', 'image/gif', 'image/webp', 'application/pdf'}
    if file.content_type not in allowed:
        raise HTTPException(status_code=400, detail="허용되지 않는 파일 형식입니다")

    file_data = await file.read()
    if len(file_data) > 10 * 1024 * 1024:  # 10MB
        raise HTTPException(status_code=400, detail="파일 크기가 10MB를 초과합니다")

    ext = file.filename.rsplit('.', 1)[-1] if '.' in file.filename else 'bin'
    filename = f"{int(datetime.utcnow().timestamp())}_{uuid.uuid4().hex[:8]}.{ext}"
    url = upload_to_local(file_data, filename, category)
    return {"url": url, "filename": filename}

# ==================== 서류 관리 API ====================

@router.get("/my/docs")
async def get_my_documents(user: dict = Depends(get_current_user)):
    """내 서류 목록 조회"""
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            SELECT id, file_category, file_name, file_path, uploaded_at as created_at
            FROM kwv_file_uploads WHERE user_id = %s
            ORDER BY uploaded_at DESC
        """, (int(user['sub']),))
        return cursor.fetchall()
    finally:
        conn.close()

@router.post("/my/docs/upload")
async def upload_my_document(
    file: UploadFile = File(...),
    category: str = Form("general"),
    user: dict = Depends(get_current_user)
):
    """서류 업로드 (인증 필요)"""
    allowed = {'image/jpeg', 'image/png', 'image/gif', 'image/webp', 'application/pdf'}
    if file.content_type not in allowed:
        raise HTTPException(status_code=400, detail="허용되지 않는 파일 형식입니다 (JPG, PNG, PDF만 가능)")

    file_data = await file.read()
    if len(file_data) > 10 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="파일 크기가 10MB를 초과합니다")

    ext = file.filename.rsplit('.', 1)[-1] if '.' in file.filename else 'bin'
    filename = f"{int(datetime.utcnow().timestamp())}_{uuid.uuid4().hex[:8]}.{ext}"
    url = upload_to_local(file_data, filename, category)

    conn = get_kwv_db_connection()
    if conn:
        try:
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO kwv_file_uploads (user_id, file_category, file_name, file_path)
                VALUES (%s, %s, %s, %s)
            """, (int(user['sub']), category, file.filename, url))
            conn.commit()
            doc_id = cursor.lastrowid
        finally:
            conn.close()
    else:
        doc_id = 0

    return {"id": doc_id, "url": url, "filename": file.filename, "category": category}

@router.delete("/my/docs/{doc_id}")
async def delete_my_document(doc_id: int, user: dict = Depends(get_current_user)):
    """내 서류 삭제"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB 연결 실패")
    try:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM kwv_file_uploads WHERE id = %s AND user_id = %s", (doc_id, int(user['sub'])))
        conn.commit()
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="서류를 찾을 수 없습니다")
        return {"message": "삭제되었습니다"}
    finally:
        conn.close()

# ==================== Google OAuth 서버 플로우 ====================

@router.get("/auth/google/start")
async def google_auth_start(request: Request, mode: str = "register"):
    """Google OAuth 시작 - Google 로그인 페이지로 리다이렉트"""
    if not GOOGLE_CLIENT_ID:
        raise HTTPException(status_code=503, detail="Google OAuth가 설정되지 않았습니다")

    # 콜백 URL 구성 (SITE_URL 우선, 없으면 요청 헤더에서 추출)
    if SITE_URL:
        base = SITE_URL.rstrip("/")
    else:
        scheme = request.headers.get("x-forwarded-proto", request.url.scheme)
        host = request.headers.get("host", request.url.hostname)
        base = f"{scheme}://{host}"
    redirect_uri = f"{base}/kwv-google-callback.html"

    google_auth_url = (
        "https://accounts.google.com/o/oauth2/v2/auth?"
        f"client_id={GOOGLE_CLIENT_ID}"
        f"&redirect_uri={redirect_uri}"
        "&response_type=code"
        "&scope=openid+profile+email"
        "&access_type=online"
        "&prompt=select_account"
        f"&state={mode}"
    )
    return RedirectResponse(url=google_auth_url)

@router.get("/auth/google/exchange")
async def google_auth_exchange(request: Request, code: str = None):
    """Google OAuth code를 사용자 정보로 교환"""
    if not code:
        raise HTTPException(status_code=400, detail="Authorization code가 필요합니다")

    if SITE_URL:
        base = SITE_URL.rstrip("/")
    else:
        scheme = request.headers.get("x-forwarded-proto", request.url.scheme)
        host = request.headers.get("host", request.url.hostname)
        base = f"{scheme}://{host}"
    redirect_uri = f"{base}/kwv-google-callback.html"

    try:
        async with httpx.AsyncClient() as client:
            # code → access_token 교환
            token_resp = await client.post("https://oauth2.googleapis.com/token", data={
                "code": code,
                "client_id": GOOGLE_CLIENT_ID,
                "client_secret": GOOGLE_CLIENT_SECRET,
                "redirect_uri": redirect_uri,
                "grant_type": "authorization_code"
            })
            token_data = token_resp.json()

            if "error" in token_data:
                raise HTTPException(status_code=400, detail=token_data.get("error_description", "Token exchange failed"))

            # access_token으로 사용자 정보 가져오기
            userinfo_resp = await client.get(
                "https://www.googleapis.com/oauth2/v3/userinfo",
                headers={"Authorization": f"Bearer {token_data['access_token']}"}
            )
            userinfo = userinfo_resp.json()

            return {
                "name": userinfo.get("name", ""),
                "email": userinfo.get("email", ""),
                "picture": userinfo.get("picture", "")
            }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Google 인증 처리 실패: {str(e)}")

@router.post("/auth/google-login")
async def google_login_by_email(request_data: GoogleLoginRequest):
    """Google 이메일로 기존 사용자 찾아 로그인"""
    email = request_data.email
    if not email:
        raise HTTPException(status_code=400, detail="이메일이 필요합니다")

    if MOCK_MODE:
        user = None
        for u in MOCK_USERS.values():
            if u["email"] == email:
                user = u
                break
        if not user:
            raise HTTPException(status_code=404, detail="등록되지 않은 사용자입니다. 먼저 회원가입을 해주세요.")

        token_data = {
            "sub": str(user["id"]),
            "email": user["email"],
            "name": user["name"],
            "user_type": user["user_type"],
            "admin_level": user["admin_level"]
        }
        access_token = create_access_token(token_data)
        return {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user["id"],
                "email": user["email"],
                "name": user["name"],
                "user_type": user["user_type"],
                "admin_level": user["admin_level"],
                "language": user.get("language", "en")
            }
        }

    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, email, name, user_type, admin_level, language
            FROM kwv_users WHERE email = %s AND is_active = TRUE
        """, (email,))
        user = cursor.fetchone()

        if not user:
            raise HTTPException(status_code=404, detail="등록되지 않은 사용자입니다. 먼저 회원가입을 해주세요.")

        user_id, email, name, user_type, admin_level, language = user
        user_type = user_type or 'applicant'
        admin_level = admin_level or 0

        cursor.execute("UPDATE kwv_users SET last_login_at = NOW(), oauth_provider = 'google' WHERE id = %s", (user_id,))
        conn.commit()

        token_data = {
            "sub": str(user_id),
            "email": email,
            "name": name,
            "user_type": user_type,
            "admin_level": admin_level
        }
        access_token = create_access_token(token_data)
        return {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user_id,
                "email": email,
                "name": name,
                "user_type": user_type,
                "admin_level": admin_level,
                "language": language or "en"
            }
        }
    finally:
        conn.close()

@router.post("/auth/admin-login")
async def admin_login(request_data: AdminLoginRequest):
    """관리자 로그인 (이름 + 비밀번호, 디폴트 비밀번호=생년월일)"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()

        # 이름으로 검색
        cursor.execute("""
            SELECT id, email, name, password_hash, password_salt, user_type, admin_level, language, is_approved
            FROM kwv_users
            WHERE name = %s AND user_type = 'admin' AND is_active = 1
        """, (request_data.name.strip(),))
        users = cursor.fetchall()

        if not users:
            raise HTTPException(status_code=404, detail="등록되지 않은 관리자입니다")

        # 비밀번호로 매칭 (동명이인도 비밀번호로 구분)
        matched_user = None
        for u in users:
            uid, email, name, pw_hash, pw_salt, user_type, admin_level, language, is_approved = u
            if pw_hash and pw_salt:
                check_hash = hashlib.sha256((request_data.password + pw_salt).encode()).hexdigest()
                if check_hash == pw_hash:
                    matched_user = u
                    break

        if not matched_user:
            raise HTTPException(status_code=401, detail="이름 또는 비밀번호가 올바르지 않습니다")

        user_id, email, name, _, _, user_type, admin_level, language, is_approved = matched_user

        if not is_approved:
            raise HTTPException(status_code=403, detail="관리자 승인 대기 중입니다. 승인 후 로그인할 수 있습니다.")

        cursor.execute("UPDATE kwv_users SET last_login_at = NOW() WHERE id = %s", (user_id,))
        conn.commit()

        token_data = {
            "sub": str(user_id),
            "email": email,
            "name": name,
            "user_type": user_type,
            "admin_level": admin_level
        }
        access_token = create_access_token(token_data)

        return {
            "access_token": access_token,
            "token_type": "bearer",
            "expires_in": JWT_EXPIRATION_HOURS * 3600,
            "user": {
                "id": user_id,
                "email": email,
                "name": name,
                "user_type": user_type,
                "admin_level": admin_level,
                "language": language or "ko"
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"로그인 실패: {str(e)}")
    finally:
        conn.close()

@router.post("/auth/admin-register")
async def admin_register(request_data: AdminRegisterRequest):
    """관리자 가입 (이름+생년월일+전화번호+소속, 비밀번호=생년월일 디폴트)"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")

    try:
        cursor = conn.cursor()

        # 중복 체크: 이름 + 생년월일 + 전화번호
        cursor.execute("""
            SELECT id FROM kwv_users
            WHERE name = %s AND birth_date = %s AND phone = %s
        """, (request_data.name.strip(), request_data.birth_date, request_data.phone.strip()))
        if cursor.fetchone():
            raise HTTPException(status_code=409, detail="이미 등록된 관리자입니다")

        # 비밀번호 = 생년월일 (디폴트)
        default_password = request_data.birth_date  # e.g. "1990-01-15"
        password_salt = secrets.token_hex(16)
        password_hash = hashlib.sha256((default_password + password_salt).encode()).hexdigest()

        # 고유 이메일 생성 (unique 제약조건 충족)
        import time
        unique_email = f"admin_{int(time.time())}_{secrets.token_hex(4)}@kwv.local"

        cursor.execute("""
            INSERT INTO kwv_users (email, name, phone, birth_date, password_hash, password_salt,
                                   user_type, admin_level, organization, region,
                                   is_approved, is_active, language)
            VALUES (%s, %s, %s, %s, %s, %s, 'admin', 2, %s, %s, FALSE, TRUE, 'ko')
        """, (unique_email, request_data.name.strip(), request_data.phone.strip(),
              request_data.birth_date, password_hash, password_salt,
              request_data.organization.strip(), request_data.organization.strip()))
        conn.commit()

        return {
            "message": "관리자 가입 신청이 완료되었습니다. 승인 후 로그인할 수 있습니다.",
            "status": "pending_approval",
            "default_password_hint": "초기 비밀번호는 생년월일입니다 (예: 1990-01-15)"
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"가입 실패: {str(e)}")
    finally:
        conn.close()

# ==================== 공개 게시판 API ====================

@router.get("/public/board")
async def get_public_board():
    """공개 구인/구직 게시판 (랜딩 페이지용, 인증 불필요)"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"jobs": [], "seekers": []}

    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)

        # 구인: 최근 활성 구인 공고 10건
        cursor.execute("""
            SELECT j.id, j.title, j.location, j.salary, j.positions, j.visa_types, j.period,
                   lg.name as lg_name, j.created_at
            FROM kwv_jobs j
            LEFT JOIN kwv_local_governments lg ON j.local_government_id = lg.id
            WHERE j.status = 'active'
            ORDER BY j.created_at DESC LIMIT 10
        """)
        jobs = cursor.fetchall()
        for j in jobs:
            if j.get('created_at'):
                j['created_at'] = j['created_at'].strftime('%Y-%m-%d') if hasattr(j['created_at'], 'strftime') else str(j['created_at'])[:10]

        # 구직: 최근 가입한 근로자(applicant) 10건 (이름 마스킹)
        cursor.execute("""
            SELECT u.id, u.name, va.nationality, va.visa_type, u.language, u.created_at
            FROM kwv_users u
            LEFT JOIN kwv_visa_applicants va ON u.id = va.user_id
            WHERE u.user_type = 'applicant' AND u.is_active = 1
            ORDER BY u.created_at DESC LIMIT 10
        """)
        seekers = cursor.fetchall()
        for s in seekers:
            # 이름 마스킹: 홍길동 → 홍*동, AB → A*
            name = s.get('name', '')
            if len(name) >= 3:
                s['name'] = name[0] + '*' * (len(name) - 2) + name[-1]
            elif len(name) == 2:
                s['name'] = name[0] + '*'
            if s.get('created_at'):
                s['created_at'] = s['created_at'].strftime('%Y-%m-%d') if hasattr(s['created_at'], 'strftime') else str(s['created_at'])[:10]

        return {"jobs": jobs, "seekers": seekers}
    except Exception as e:
        print(f"public/board error: {e}")
        return {"jobs": [], "seekers": []}
    finally:
        conn.close()

# ==================== 시스템 설정 API ====================

def ensure_system_settings_table():
    """시스템 설정 테이블 확인 및 생성"""
    conn = get_kwv_db_connection()
    if not conn:
        return
    try:
        cursor = conn.cursor()
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS kwv_system_settings (
                id INT AUTO_INCREMENT PRIMARY KEY,
                setting_key VARCHAR(100) UNIQUE NOT NULL,
                setting_value TEXT,
                setting_type ENUM('string','number','boolean','json') DEFAULT 'string',
                description VARCHAR(500),
                updated_by INT,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP ON UPDATE CURRENT_TIMESTAMP,
                INDEX idx_key (setting_key)
            ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4 COLLATE=utf8mb4_unicode_ci
        """)
        conn.commit()
    except:
        pass
    finally:
        conn.close()

@router.get("/dashboard/stats")
async def get_dashboard_stats_real():
    """대시보드 통계 (실제 DB 조회)"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"success": True, "data": {"total_jobs": 0, "total_positions": 0, "allocated_quota": 0, "kr_regions": 0, "foreign_regions": 0, "visa_types": 0}}
    try:
        cursor = conn.cursor()
        # 구인 통계
        cursor.execute("SELECT COUNT(*), COALESCE(SUM(positions),0) FROM kwv_jobs WHERE status='active'")
        jobs_row = cursor.fetchone()
        total_jobs = jobs_row[0] or 0
        total_positions = int(jobs_row[1] or 0)
        # 모집 가능 TO (활성 지자체 allocated_quota 합산)
        cursor.execute("SELECT COALESCE(SUM(allocated_quota),0) FROM kwv_local_governments WHERE is_active=1")
        allocated_quota = int(cursor.fetchone()[0] or 0)
        # 한국 지자체 (활성만)
        cursor.execute("SELECT COUNT(*) FROM kwv_local_governments WHERE is_active=1")
        kr_regions = cursor.fetchone()[0] or 0
        # 해외 기관
        cursor.execute("SELECT COUNT(DISTINCT partner_country) FROM kwv_mou_agreements")
        foreign_regions = cursor.fetchone()[0] or 0
        # 비자 종류 (신청자 기준)
        cursor.execute("SELECT COUNT(DISTINCT visa_type) FROM kwv_visa_applicants WHERE visa_type IS NOT NULL AND visa_type != ''")
        visa_types = cursor.fetchone()[0] or 0
        return {"success": True, "data": {
            "total_jobs": total_jobs, "total_positions": total_positions,
            "allocated_quota": allocated_quota,
            "kr_regions": kr_regions, "foreign_regions": foreign_regions, "visa_types": visa_types
        }}
    finally:
        conn.close()

@router.post("/admin/logo/upload")
async def upload_logo(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    """로고 이미지 업로드"""
    require_admin_level(user, 9)
    allowed = {'image/jpeg', 'image/png', 'image/gif', 'image/webp', 'image/svg+xml'}
    if file.content_type not in allowed:
        raise HTTPException(status_code=400, detail="이미지 파일만 업로드 가능합니다")
    file_data = await file.read()
    if len(file_data) > 5 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="파일 크기가 5MB를 초과합니다")
    ext = file.filename.rsplit('.', 1)[-1] if '.' in file.filename else 'png'
    filename = f"logo_{int(datetime.utcnow().timestamp())}.{ext}"
    url = upload_to_local(file_data, filename, "logos")
    # 시스템 설정에 로고 URL 저장
    conn = get_kwv_db_connection()
    if conn:
        try:
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO kwv_system_settings (setting_key, setting_value, updated_by)
                VALUES ('logo_url', %s, %s)
                ON DUPLICATE KEY UPDATE setting_value = VALUES(setting_value), updated_by = VALUES(updated_by)
            """, (url, user.get('sub')))
            conn.commit()
        finally:
            conn.close()
    return {"url": url, "filename": filename}

@router.post("/admin/test-api-key")
async def test_api_key(request: Request, user: dict = Depends(get_current_user)):
    """API 키 테스트"""
    require_admin_level(user, 9)
    body = await request.json()
    key_type = body.get("type", "")
    api_key = body.get("key", "")

    if key_type == "google_maps":
        try:
            async with httpx.AsyncClient() as client:
                res = await client.get(f"https://maps.googleapis.com/maps/api/geocode/json?address=Seoul&key={api_key}", timeout=5)
                data = res.json()
                if data.get("status") == "OK":
                    return {"success": True, "message": "Google Maps API 키가 유효합니다"}
                elif data.get("status") == "REQUEST_DENIED":
                    return {"success": False, "message": f"API 키가 거부되었습니다: {data.get('error_message','')}"}
                else:
                    return {"success": False, "message": f"응답: {data.get('status','')}"}
        except Exception as e:
            return {"success": False, "message": f"연결 오류: {str(e)}"}
    elif key_type == "ai":
        if not api_key:
            return {"success": False, "message": "API 키를 입력하세요"}
        # Groq API 테스트
        try:
            async with httpx.AsyncClient() as client:
                res = await client.get("https://api.groq.com/openai/v1/models",
                    headers={"Authorization": f"Bearer {api_key}"}, timeout=5)
                if res.status_code == 200:
                    return {"success": True, "message": "AI API 키가 유효합니다 (Groq)"}
                else:
                    return {"success": False, "message": f"Groq 응답: {res.status_code}"}
        except:
            pass
        # Gemini API 테스트
        try:
            async with httpx.AsyncClient() as client:
                res = await client.get(f"https://generativelanguage.googleapis.com/v1beta/models?key={api_key}", timeout=5)
                if res.status_code == 200:
                    return {"success": True, "message": "AI API 키가 유효합니다 (Gemini)"}
                else:
                    return {"success": False, "message": f"API 키를 확인해주세요"}
        except Exception as e:
            return {"success": False, "message": f"연결 오류: {str(e)}"}
    else:
        return {"success": False, "message": "알 수 없는 키 타입"}

@router.get("/settings")
async def get_system_settings():
    """시스템 설정 조회 (공개)"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT setting_key, setting_value, setting_type FROM kwv_system_settings")
        rows = cursor.fetchall()
        settings = {}
        for key, value, stype in rows:
            if stype == 'boolean':
                settings[key] = value == 'true'
            elif stype == 'number':
                try:
                    settings[key] = int(value) if value else 0
                except:
                    settings[key] = float(value) if value else 0
            elif stype == 'json':
                try:
                    settings[key] = json.loads(value) if value else {}
                except:
                    settings[key] = value
            else:
                settings[key] = value or ''
        # 기본 로고 URL 폴백
        if not settings.get('logo_url'):
            settings['logo_url'] = '/woosong-logo.png'
        return settings
    finally:
        conn.close()

@router.post("/settings")
async def update_system_settings(request: Request, user: dict = Depends(get_current_user)):
    """시스템 설정 수정 (super admin 전용)"""
    require_admin_level(user, 9)
    body = await request.json()

    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        updated = 0
        for key, value in body.items():
            if isinstance(value, bool):
                str_value = 'true' if value else 'false'
            elif isinstance(value, (dict, list)):
                str_value = json.dumps(value)
            else:
                str_value = str(value)
            cursor.execute("""
                INSERT INTO kwv_system_settings (setting_key, setting_value, updated_by)
                VALUES (%s, %s, %s)
                ON DUPLICATE KEY UPDATE setting_value = VALUES(setting_value), updated_by = VALUES(updated_by)
            """, (key, str_value, user.get('sub')))
            updated += 1
        conn.commit()
        return {"message": f"{updated}개 설정이 저장되었습니다", "updated": updated}
    finally:
        conn.close()

# ==================== 지자체 API ====================

class LocalGovernmentCreate(BaseModel):
    name: str
    name_en: Optional[str] = None
    region: str
    address: Optional[str] = None
    phone: Optional[str] = None
    email: Optional[str] = None
    website_url: Optional[str] = None
    representative_name: Optional[str] = None
    representative_phone: Optional[str] = None
    representative_email: Optional[str] = None
    allocated_quota: Optional[int] = 0
    quota_year: Optional[int] = None
    logo_url: Optional[str] = None
    description: Optional[str] = None
    description_en: Optional[str] = None
    latitude: Optional[float] = None
    longitude: Optional[float] = None

@router.get("/local-governments")
async def list_local_governments(region: Optional[str] = None, active_only: bool = True):
    """지자체 목록 조회 (공개)"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        query = """
            SELECT id, name, name_en, region, address, phone, email, website_url,
                   representative_name, representative_phone, representative_email,
                   allocated_quota, used_quota, quota_year,
                   logo_url, description, description_en, latitude, longitude,
                   is_active, created_at
            FROM kwv_local_governments
        """
        conditions = []
        params = []
        if active_only:
            conditions.append("is_active = TRUE")
        if region:
            conditions.append("region = %s")
            params.append(region)
        if conditions:
            query += " WHERE " + " AND ".join(conditions)
        query += " ORDER BY region, name"
        cursor.execute(query, params)
        rows = cursor.fetchall()
        columns = [desc[0] for desc in cursor.description]
        result = []
        for row in rows:
            item = dict(zip(columns, row))
            # datetime/decimal 변환
            for k, v in item.items():
                if isinstance(v, datetime):
                    item[k] = v.isoformat()
                elif hasattr(v, '__float__'):
                    item[k] = float(v)
            result.append(item)
        return result
    finally:
        conn.close()

@router.get("/local-governments/{lg_id}")
async def get_local_government(lg_id: int):
    """지자체 상세 조회"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, name, name_en, region, address, phone, email, website_url,
                   representative_name, representative_phone, representative_email,
                   allocated_quota, used_quota, quota_year,
                   logo_url, description, description_en, latitude, longitude,
                   is_active, created_at, updated_at
            FROM kwv_local_governments WHERE id = %s
        """, (lg_id,))
        row = cursor.fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="지자체를 찾을 수 없습니다")
        columns = [desc[0] for desc in cursor.description]
        item = dict(zip(columns, row))
        for k, v in item.items():
            if isinstance(v, datetime):
                item[k] = v.isoformat()
            elif hasattr(v, '__float__'):
                item[k] = float(v)
        # 배정된 근로자 수 조회
        cursor.execute("SELECT COUNT(*) FROM kwv_users WHERE local_government_id = %s AND is_active = TRUE", (lg_id,))
        item['worker_count'] = cursor.fetchone()[0]
        return item
    finally:
        conn.close()

@router.post("/local-governments")
async def create_local_government(data: LocalGovernmentCreate, user: dict = Depends(get_current_user)):
    """지자체 등록 (admin)"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        from datetime import date
        cursor.execute("""
            INSERT INTO kwv_local_governments
            (name, name_en, region, address, phone, email, website_url,
             representative_name, representative_phone, representative_email,
             allocated_quota, quota_year, logo_url, description, description_en,
             latitude, longitude)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
        """, (data.name, data.name_en, data.region, data.address, data.phone,
              data.email, data.website_url, data.representative_name,
              data.representative_phone, data.representative_email,
              data.allocated_quota, data.quota_year or date.today().year,
              data.logo_url, data.description, data.description_en,
              data.latitude, data.longitude))
        conn.commit()
        return {"id": cursor.lastrowid, "message": f"지자체 '{data.name}' 등록 완료"}
    finally:
        conn.close()

@router.put("/local-governments/{lg_id}")
async def update_local_government(lg_id: int, request: Request, user: dict = Depends(get_current_user)):
    """지자체 수정"""
    require_admin(user)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        allowed_fields = [
            'name', 'name_en', 'region', 'address', 'phone', 'email', 'website_url',
            'representative_name', 'representative_phone', 'representative_email',
            'allocated_quota', 'used_quota', 'quota_year', 'logo_url',
            'description', 'description_en', 'latitude', 'longitude', 'is_active'
        ]
        updates = []
        params = []
        for field in allowed_fields:
            if field in body:
                updates.append(f"{field} = %s")
                params.append(body[field])
        if not updates:
            raise HTTPException(status_code=400, detail="수정할 항목이 없습니다")
        params.append(lg_id)
        cursor.execute(f"UPDATE kwv_local_governments SET {', '.join(updates)} WHERE id = %s", params)
        conn.commit()
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="지자체를 찾을 수 없습니다")
        return {"message": "지자체 정보가 수정되었습니다"}
    finally:
        conn.close()

@router.delete("/local-governments/{lg_id}")
async def delete_local_government(lg_id: int, user: dict = Depends(get_current_user)):
    """지자체 삭제 (super admin)"""
    require_admin_level(user, 9)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("UPDATE kwv_local_governments SET is_active = FALSE WHERE id = %s", (lg_id,))
        conn.commit()
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="지자체를 찾을 수 없습니다")
        return {"message": "지자체가 비활성화되었습니다"}
    finally:
        conn.close()

@router.put("/local-governments/{lg_id}/quota")
async def update_quota(lg_id: int, request: Request, user: dict = Depends(get_current_user)):
    """지자체 TO 배정 관리 (super admin)"""
    require_admin_level(user, 9)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        allocated = body.get('allocated_quota')
        year = body.get('quota_year')
        if allocated is not None:
            cursor.execute("""
                UPDATE kwv_local_governments
                SET allocated_quota = %s, quota_year = COALESCE(%s, quota_year)
                WHERE id = %s
            """, (allocated, year, lg_id))
        conn.commit()
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="지자체를 찾을 수 없습니다")
        return {"message": "TO 배정이 업데이트되었습니다"}
    finally:
        conn.close()

# ==================== MOU 협정 API ====================

class MouCreate(BaseModel):
    title: str
    title_en: Optional[str] = None
    partner_country: str
    partner_country_name: Optional[str] = None
    partner_type: Optional[str] = "government"
    partner_organization: str
    partner_organization_en: Optional[str] = None
    partner_representative: Optional[str] = None
    partner_contact: Optional[str] = None
    korean_organization: Optional[str] = None
    korean_representative: Optional[str] = None
    description: Optional[str] = None
    description_en: Optional[str] = None
    signed_date: Optional[str] = None
    effective_date: Optional[str] = None
    expiry_date: Optional[str] = None
    worker_quota: Optional[int] = 0
    document_url: Optional[str] = None
    photo_url: Optional[str] = None
    photo_url_2: Optional[str] = None
    photo_url_3: Optional[str] = None
    status: Optional[str] = "draft"
    is_public: Optional[bool] = False
    display_order: Optional[int] = 0

@router.get("/mou")
async def list_mou(
    country: Optional[str] = None,
    status: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """MOU 목록 (관리자)"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        where = "WHERE 1=1"
        params = []
        if country:
            where += " AND partner_country = %s"
            params.append(country)
        if status:
            where += " AND status = %s"
            params.append(status)
        cursor.execute(f"""
            SELECT id, title, title_en, partner_country, partner_country_name, partner_type,
                   partner_organization, signed_date, expiry_date, worker_quota,
                   photo_url, status, is_public, display_order, created_at
            FROM kwv_mou_agreements {where} ORDER BY display_order ASC, signed_date DESC
        """, params)
        mous = []
        for r in cursor.fetchall():
            mous.append({
                "id": r[0], "title": r[1], "title_en": r[2],
                "partner_country": r[3], "partner_country_name": r[4],
                "partner_type": r[5], "partner_organization": r[6],
                "signed_date": r[7].isoformat() if r[7] else None,
                "expiry_date": r[8].isoformat() if r[8] else None,
                "worker_quota": r[9], "photo_url": r[10],
                "status": r[11], "is_public": bool(r[12]),
                "display_order": r[13],
                "created_at": r[14].isoformat() if r[14] else None
            })
        return mous
    finally:
        conn.close()

@router.get("/mou/public")
async def list_public_mou():
    """공개 MOU 목록 (로그인 불요)"""
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, title, title_en, partner_country, partner_country_name, partner_type,
                   partner_organization, partner_organization_en, partner_representative,
                   korean_organization, korean_representative,
                   description, description_en, signed_date, expiry_date, worker_quota,
                   photo_url, photo_url_2, photo_url_3, status, display_order
            FROM kwv_mou_agreements WHERE is_public = TRUE AND status IN ('active','draft')
            ORDER BY display_order ASC, signed_date DESC
        """)
        mous = []
        for r in cursor.fetchall():
            mous.append({
                "id": r[0], "title": r[1], "title_en": r[2],
                "partner_country": r[3], "partner_country_name": r[4],
                "partner_type": r[5], "partner_organization": r[6],
                "partner_organization_en": r[7], "partner_representative": r[8],
                "korean_organization": r[9], "korean_representative": r[10],
                "description": r[11], "description_en": r[12],
                "signed_date": r[13].isoformat() if r[13] else None,
                "expiry_date": r[14].isoformat() if r[14] else None,
                "worker_quota": r[15], "photo_url": r[16],
                "photo_url_2": r[17], "photo_url_3": r[18],
                "status": r[19], "display_order": r[20]
            })
        return mous
    finally:
        conn.close()

@router.get("/mou/{mou_id}")
async def get_mou(mou_id: int, user: dict = Depends(get_current_user)):
    """MOU 상세"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM kwv_mou_agreements WHERE id = %s", (mou_id,))
        r = cursor.fetchone()
        if not r:
            raise HTTPException(status_code=404, detail="MOU를 찾을 수 없습니다")
        cols = [d[0] for d in cursor.description]
        mou = {}
        for i, col in enumerate(cols):
            val = r[i]
            if hasattr(val, 'isoformat'):
                val = val.isoformat()
            if col == 'is_public':
                val = bool(val)
            mou[col] = val
        return mou
    finally:
        conn.close()

@router.post("/mou")
async def create_mou(mou_data: MouCreate, user: dict = Depends(get_current_user)):
    """MOU 등록"""
    require_admin_level(user, 2)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO kwv_mou_agreements (
                title, title_en, partner_country, partner_country_name, partner_type,
                partner_organization, partner_organization_en, partner_representative, partner_contact,
                korean_organization, korean_representative,
                description, description_en, signed_date, effective_date, expiry_date, worker_quota,
                document_url, photo_url, photo_url_2, photo_url_3,
                status, is_public, display_order, created_by
            ) VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
        """, (
            mou_data.title, mou_data.title_en, mou_data.partner_country, mou_data.partner_country_name,
            mou_data.partner_type, mou_data.partner_organization, mou_data.partner_organization_en,
            mou_data.partner_representative, mou_data.partner_contact,
            mou_data.korean_organization, mou_data.korean_representative,
            mou_data.description, mou_data.description_en,
            mou_data.signed_date or None, mou_data.effective_date or None, mou_data.expiry_date or None,
            mou_data.worker_quota, mou_data.document_url, mou_data.photo_url,
            mou_data.photo_url_2, mou_data.photo_url_3,
            mou_data.status, mou_data.is_public, mou_data.display_order, user.get('sub')
        ))
        conn.commit()
        return {"id": cursor.lastrowid, "message": f"MOU '{mou_data.title}' 등록 완료"}
    finally:
        conn.close()

@router.put("/mou/{mou_id}")
async def update_mou(mou_id: int, request: Request, user: dict = Depends(get_current_user)):
    """MOU 수정"""
    require_admin_level(user, 2)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        allowed = ['title','title_en','partner_country','partner_country_name','partner_type',
                    'partner_organization','partner_organization_en','partner_representative','partner_contact',
                    'korean_organization','korean_representative','description','description_en',
                    'signed_date','effective_date','expiry_date','worker_quota',
                    'document_url','photo_url','photo_url_2','photo_url_3',
                    'status','is_public','display_order']
        sets = []
        params = []
        for k, v in body.items():
            if k in allowed:
                sets.append(f"{k} = %s")
                params.append(v)
        if not sets:
            raise HTTPException(status_code=400, detail="수정할 항목이 없습니다")
        params.append(mou_id)
        cursor.execute(f"UPDATE kwv_mou_agreements SET {', '.join(sets)} WHERE id = %s", params)
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="MOU를 찾을 수 없습니다")
        conn.commit()
        return {"message": "MOU가 수정되었습니다"}
    finally:
        conn.close()

@router.delete("/mou/{mou_id}")
async def delete_mou(mou_id: int, user: dict = Depends(get_current_user)):
    """MOU 삭제"""
    require_admin_level(user, 9)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="Database connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM kwv_mou_agreements WHERE id = %s", (mou_id,))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="MOU를 찾을 수 없습니다")
        conn.commit()
        return {"message": "MOU가 삭제되었습니다"}
    finally:
        conn.close()

# ==================== 통계 API 강화 ====================

@router.get("/admin/statistics/by-nationality")
async def stats_by_nationality(user: dict = Depends(get_current_user)):
    """국적별 근로자 통계"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT a.nationality, COUNT(*) as cnt
            FROM kwv_visa_applicants a
            JOIN kwv_users u ON a.user_id = u.id
            WHERE u.user_type = 'applicant' AND a.nationality IS NOT NULL
            GROUP BY a.nationality ORDER BY cnt DESC
        """)
        return [{"nationality": r[0], "count": r[1]} for r in cursor.fetchall()]
    finally:
        conn.close()

@router.get("/admin/statistics/by-region")
async def stats_by_region(user: dict = Depends(get_current_user)):
    """지역별 근로자 분포"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT lg.region, lg.name, COUNT(u.id) as worker_count, lg.allocated_quota, lg.used_quota
            FROM kwv_local_governments lg
            LEFT JOIN kwv_users u ON u.local_government_id = lg.id AND u.user_type = 'applicant'
            WHERE lg.is_active = TRUE
            GROUP BY lg.id ORDER BY worker_count DESC
        """)
        return [{"region": r[0], "lg_name": r[1], "worker_count": r[2],
                 "allocated_quota": r[3], "used_quota": r[4]} for r in cursor.fetchall()]
    finally:
        conn.close()

@router.get("/admin/statistics/by-visa")
async def stats_by_visa(user: dict = Depends(get_current_user)):
    """비자유형별 통계"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT a.visa_type, COUNT(*) as cnt
            FROM kwv_visa_applicants a
            JOIN kwv_users u ON a.user_id = u.id
            WHERE u.user_type = 'applicant' AND a.visa_type IS NOT NULL
            GROUP BY a.visa_type ORDER BY cnt DESC
        """)
        return [{"visa_type": r[0], "count": r[1]} for r in cursor.fetchall()]
    finally:
        conn.close()

@router.get("/admin/statistics/monthly")
async def stats_monthly(user: dict = Depends(get_current_user)):
    """월별 가입 추이"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT DATE_FORMAT(created_at, '%%Y-%%m') as month, COUNT(*) as cnt
            FROM kwv_users WHERE user_type = 'applicant'
            GROUP BY month ORDER BY month DESC LIMIT 12
        """)
        return [{"month": r[0], "count": r[1]} for r in cursor.fetchall()]
    finally:
        conn.close()

# ==================== Phase 4: 출퇴근 시스템 ====================

class WorkplaceCreate(BaseModel):
    local_government_id: int
    name: str
    name_en: Optional[str] = None
    address: Optional[str] = None
    latitude: Optional[float] = None
    longitude: Optional[float] = None
    geofence_radius: int = 200
    manager_name: Optional[str] = None
    manager_phone: Optional[str] = None
    worker_capacity: int = 0

class AttendanceCheck(BaseModel):
    workplace_id: Optional[int] = None
    qr_code: Optional[str] = None
    check_type: str  # check_in / check_out
    latitude: Optional[float] = None
    longitude: Optional[float] = None
    photo_url: Optional[str] = None
    note: Optional[str] = None

import math

def haversine(lat1, lon1, lat2, lon2):
    """두 좌표 간 거리(미터) 계산"""
    R = 6371000
    phi1, phi2 = math.radians(lat1), math.radians(lat2)
    dphi = math.radians(lat2 - lat1)
    dlambda = math.radians(lon2 - lon1)
    a = math.sin(dphi/2)**2 + math.cos(phi1)*math.cos(phi2)*math.sin(dlambda/2)**2
    return R * 2 * math.atan2(math.sqrt(a), math.sqrt(1-a))

# --- 사업장 CRUD ---

@router.get("/workplaces")
async def list_workplaces(
    lg_id: Optional[int] = None,
    user: dict = Depends(get_current_user)
):
    """사업장 목록"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        where = "WHERE 1=1"
        params = []
        if lg_id:
            where += " AND w.local_government_id = %s"
            params.append(lg_id)
        cursor.execute(f"""
            SELECT w.id, w.name, w.name_en, w.address, w.latitude, w.longitude,
                   w.geofence_radius, w.qr_code, w.manager_name, w.manager_phone,
                   w.worker_capacity, w.is_active, w.local_government_id,
                   lg.name as lg_name,
                   (SELECT COUNT(*) FROM kwv_worker_assignments wa WHERE wa.workplace_id = w.id AND wa.status = 'active') as worker_count
            FROM kwv_workplaces w
            LEFT JOIN kwv_local_governments lg ON w.local_government_id = lg.id
            {where} ORDER BY w.name
        """, params)
        result = []
        for r in cursor.fetchall():
            result.append({
                "id": r[0], "name": r[1], "name_en": r[2], "address": r[3],
                "latitude": float(r[4]) if r[4] else None,
                "longitude": float(r[5]) if r[5] else None,
                "geofence_radius": r[6], "qr_code": r[7],
                "manager_name": r[8], "manager_phone": r[9],
                "worker_capacity": r[10], "is_active": bool(r[11]),
                "local_government_id": r[12], "lg_name": r[13],
                "worker_count": r[14]
            })
        return result
    finally:
        conn.close()

@router.post("/workplaces")
async def create_workplace(data: WorkplaceCreate, user: dict = Depends(get_current_user)):
    """사업장 등록"""
    require_admin_level(user, 2)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        qr_code = f"WP-{uuid.uuid4().hex[:8].upper()}"
        cursor.execute("""
            INSERT INTO kwv_workplaces (local_government_id, name, name_en, address,
                latitude, longitude, geofence_radius, qr_code, manager_name, manager_phone, worker_capacity)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
        """, (data.local_government_id, data.name, data.name_en, data.address,
              data.latitude, data.longitude, data.geofence_radius, qr_code,
              data.manager_name, data.manager_phone, data.worker_capacity))
        conn.commit()
        return {"id": cursor.lastrowid, "qr_code": qr_code, "message": f"사업장 '{data.name}' 등록 완료"}
    finally:
        conn.close()

@router.get("/workplaces/{wp_id}")
async def get_workplace(wp_id: int, user: dict = Depends(get_current_user)):
    """사업장 상세"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM kwv_workplaces WHERE id = %s", (wp_id,))
        r = cursor.fetchone()
        if not r:
            raise HTTPException(status_code=404, detail="사업장을 찾을 수 없습니다")
        cols = [d[0] for d in cursor.description]
        wp = {}
        for i, col in enumerate(cols):
            val = r[i]
            if hasattr(val, 'isoformat'):
                val = val.isoformat()
            if isinstance(val, bytes):
                val = bool(val[0]) if len(val) == 1 else val.decode()
            if col in ('latitude', 'longitude') and val is not None:
                val = float(val)
            wp[col] = val
        return wp
    finally:
        conn.close()

@router.put("/workplaces/{wp_id}")
async def update_workplace(wp_id: int, request: Request, user: dict = Depends(get_current_user)):
    """사업장 수정"""
    require_admin_level(user, 2)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        allowed = ['name','name_en','address','latitude','longitude','geofence_radius',
                    'manager_name','manager_phone','worker_capacity','is_active','local_government_id']
        sets, params = [], []
        for k, v in body.items():
            if k in allowed:
                sets.append(f"{k} = %s")
                params.append(v)
        if not sets:
            raise HTTPException(status_code=400, detail="수정할 항목이 없습니다")
        params.append(wp_id)
        cursor.execute(f"UPDATE kwv_workplaces SET {', '.join(sets)} WHERE id = %s", params)
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="사업장을 찾을 수 없습니다")
        conn.commit()
        return {"message": "사업장이 수정되었습니다"}
    finally:
        conn.close()

@router.delete("/workplaces/{wp_id}")
async def delete_workplace(wp_id: int, user: dict = Depends(get_current_user)):
    """사업장 삭제"""
    require_admin_level(user, 9)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM kwv_workplaces WHERE id = %s", (wp_id,))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="사업장을 찾을 수 없습니다")
        conn.commit()
        return {"message": "사업장이 삭제되었습니다"}
    finally:
        conn.close()

@router.post("/workplaces/{wp_id}/regenerate-qr")
async def regenerate_qr(wp_id: int, user: dict = Depends(get_current_user)):
    """QR 코드 재생성"""
    require_admin_level(user, 2)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        new_qr = f"WP-{uuid.uuid4().hex[:8].upper()}"
        cursor.execute("UPDATE kwv_workplaces SET qr_code = %s WHERE id = %s", (new_qr, wp_id))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="사업장을 찾을 수 없습니다")
        conn.commit()
        return {"qr_code": new_qr, "message": "QR 코드가 재생성되었습니다"}
    finally:
        conn.close()

# --- 근로자-사업장 배정 ---

@router.post("/workplaces/{wp_id}/assign")
async def assign_worker(wp_id: int, request: Request, user: dict = Depends(get_current_user)):
    """근로자를 사업장에 배정"""
    require_admin_level(user, 2)
    body = await request.json()
    user_id = body.get("user_id")
    if not user_id:
        raise HTTPException(status_code=400, detail="user_id가 필요합니다")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT id FROM kwv_workplaces WHERE id = %s", (wp_id,))
        if not cursor.fetchone():
            raise HTTPException(status_code=404, detail="사업장을 찾을 수 없습니다")
        cursor.execute("""
            SELECT id FROM kwv_worker_assignments
            WHERE user_id = %s AND workplace_id = %s AND status = 'active'
        """, (user_id, wp_id))
        if cursor.fetchone():
            raise HTTPException(status_code=409, detail="이미 배정된 근로자입니다")
        from datetime import date
        cursor.execute("""
            INSERT INTO kwv_worker_assignments (user_id, workplace_id, assigned_date)
            VALUES (%s, %s, %s)
        """, (user_id, wp_id, date.today()))
        conn.commit()
        return {"message": "근로자가 사업장에 배정되었습니다"}
    finally:
        conn.close()

@router.get("/workplaces/{wp_id}/workers")
async def workplace_workers(wp_id: int, user: dict = Depends(get_current_user)):
    """사업장 소속 근로자 목록"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT wa.id, wa.user_id, u.name, u.email, u.phone, u.profile_photo,
                   va.nationality, va.visa_type, wa.assigned_date, wa.status
            FROM kwv_worker_assignments wa
            JOIN kwv_users u ON wa.user_id = u.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            WHERE wa.workplace_id = %s AND wa.status = 'active'
            ORDER BY u.name
        """, (wp_id,))
        return [{
            "assignment_id": r[0], "user_id": r[1], "name": r[2], "email": r[3],
            "phone": r[4], "profile_photo": r[5], "nationality": r[6],
            "visa_type": r[7], "assigned_date": r[8].isoformat() if r[8] else None,
            "status": r[9]
        } for r in cursor.fetchall()]
    finally:
        conn.close()

@router.delete("/worker-assignments/{assignment_id}")
async def remove_assignment(assignment_id: int, user: dict = Depends(get_current_user)):
    """근로자 사업장 배정 해제"""
    require_admin_level(user, 2)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("UPDATE kwv_worker_assignments SET status = 'completed' WHERE id = %s AND status = 'active'", (assignment_id,))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="배정을 찾을 수 없습니다")
        conn.commit()
        return {"message": "배정이 해제되었습니다"}
    finally:
        conn.close()

# --- 출퇴근 체크인/체크아웃 ---

@router.post("/attendance/check")
async def attendance_check(data: AttendanceCheck, request: Request, user: dict = Depends(get_current_user)):
    """출퇴근 체크 (QR / GPS)"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        # QR 코드로 사업장 찾기
        workplace_id = data.workplace_id
        if data.qr_code and not workplace_id:
            cursor.execute("SELECT id FROM kwv_workplaces WHERE qr_code = %s AND is_active = 1", (data.qr_code,))
            wp = cursor.fetchone()
            if not wp:
                raise HTTPException(status_code=404, detail="유효하지 않은 QR 코드입니다")
            workplace_id = wp[0]

        if not workplace_id:
            raise HTTPException(status_code=400, detail="사업장 ID 또는 QR 코드가 필요합니다")

        # 사업장 정보
        cursor.execute("SELECT latitude, longitude, geofence_radius, name FROM kwv_workplaces WHERE id = %s", (workplace_id,))
        wp_info = cursor.fetchone()
        if not wp_info:
            raise HTTPException(status_code=404, detail="사업장을 찾을 수 없습니다")
        wp_lat, wp_lon, wp_radius, wp_name = wp_info
        wp_radius = wp_radius or 200

        # GPS 거리 검증
        distance = None
        is_valid = True
        invalid_reason = None
        check_method = 'qr' if data.qr_code else 'gps'

        if data.latitude and data.longitude and wp_lat and wp_lon:
            distance = int(haversine(float(data.latitude), float(data.longitude), float(wp_lat), float(wp_lon)))
            if distance > wp_radius:
                is_valid = False
                invalid_reason = f"사업장에서 {distance}m 떨어져 있습니다 (허용: {wp_radius}m)"

        # 중복 체크 (같은 날 같은 타입 5분 이내 중복 방지)
        user_id = user.get("user_id") or user.get("sub")
        cursor.execute("""
            SELECT id FROM kwv_attendance
            WHERE user_id = %s AND workplace_id = %s AND check_type = %s
            AND check_time > DATE_SUB(NOW(), INTERVAL 5 MINUTE)
        """, (user_id, workplace_id, data.check_type))
        if cursor.fetchone():
            raise HTTPException(status_code=429, detail="5분 이내 중복 체크입니다")

        # IP 주소
        ip = request.client.host if request.client else None

        cursor.execute("""
            INSERT INTO kwv_attendance (user_id, workplace_id, check_type, check_method,
                latitude, longitude, distance_from_workplace, is_valid, invalid_reason,
                photo_url, ip_address, note)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
        """, (user_id, workplace_id, data.check_type, check_method,
              data.latitude, data.longitude, distance, is_valid, invalid_reason,
              data.photo_url, ip, data.note))
        conn.commit()

        status_text = "출근" if data.check_type == "check_in" else "퇴근"
        return {
            "id": cursor.lastrowid,
            "message": f"{wp_name} {status_text} 완료",
            "is_valid": is_valid,
            "distance": distance,
            "check_time": datetime.utcnow().isoformat(),
            "invalid_reason": invalid_reason
        }
    finally:
        conn.close()

@router.get("/attendance/my")
async def my_attendance(
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """내 출퇴근 기록"""
    user_id = user.get("user_id") or user.get("sub")
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        where = "WHERE a.user_id = %s"
        params = [user_id]
        if date_from:
            where += " AND DATE(a.check_time) >= %s"
            params.append(date_from)
        if date_to:
            where += " AND DATE(a.check_time) <= %s"
            params.append(date_to)
        cursor.execute(f"""
            SELECT a.id, a.check_type, a.check_method, a.check_time, a.is_valid,
                   a.distance_from_workplace, a.invalid_reason, w.name as workplace_name
            FROM kwv_attendance a
            JOIN kwv_workplaces w ON a.workplace_id = w.id
            {where} ORDER BY a.check_time DESC LIMIT 100
        """, params)
        return [{
            "id": r[0], "check_type": r[1], "check_method": r[2],
            "check_time": r[3].isoformat() if r[3] else None,
            "is_valid": bool(r[4]), "distance": r[5], "invalid_reason": r[6],
            "workplace_name": r[7]
        } for r in cursor.fetchall()]
    finally:
        conn.close()

@router.get("/attendance/today")
async def my_today(user: dict = Depends(get_current_user)):
    """오늘의 출퇴근 상태"""
    user_id = user.get("user_id") or user.get("sub")
    conn = get_kwv_db_connection()
    if not conn:
        return {"checked_in": False, "records": []}
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT a.id, a.check_type, a.check_time, a.is_valid, w.name
            FROM kwv_attendance a
            JOIN kwv_workplaces w ON a.workplace_id = w.id
            WHERE a.user_id = %s AND DATE(a.check_time) = CURDATE()
            ORDER BY a.check_time ASC
        """, (user_id,))
        records = []
        last_type = None
        for r in cursor.fetchall():
            records.append({
                "id": r[0], "check_type": r[1],
                "check_time": r[2].isoformat() if r[2] else None,
                "is_valid": bool(r[3]), "workplace_name": r[4]
            })
            last_type = r[1]
        return {
            "checked_in": last_type == "check_in",
            "records": records
        }
    finally:
        conn.close()

# --- 관리자 출퇴근 조회 ---

@router.get("/attendance/admin")
async def admin_attendance(
    workplace_id: Optional[int] = None,
    user_id: Optional[int] = None,
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    is_valid: Optional[str] = None,
    page: int = 1,
    per_page: int = 50,
    user: dict = Depends(get_current_user)
):
    """관리자 출퇴근 조회"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": [], "total": 0}
    try:
        cursor = conn.cursor()
        where = "WHERE 1=1"
        params = []
        if workplace_id:
            where += " AND a.workplace_id = %s"
            params.append(workplace_id)
        if user_id:
            where += " AND a.user_id = %s"
            params.append(user_id)
        if date_from:
            where += " AND DATE(a.check_time) >= %s"
            params.append(date_from)
        if date_to:
            where += " AND DATE(a.check_time) <= %s"
            params.append(date_to)
        if is_valid == 'true':
            where += " AND a.is_valid = 1"
        elif is_valid == 'false':
            where += " AND a.is_valid = 0"

        cursor.execute(f"SELECT COUNT(*) FROM kwv_attendance a {where}", params)
        total = cursor.fetchone()[0]

        offset = (page - 1) * per_page
        cursor.execute(f"""
            SELECT a.id, a.user_id, u.name as user_name, u.profile_photo,
                   va.nationality, a.workplace_id, w.name as workplace_name,
                   a.check_type, a.check_method, a.check_time, a.is_valid,
                   a.distance_from_workplace, a.invalid_reason, a.latitude, a.longitude
            FROM kwv_attendance a
            JOIN kwv_users u ON a.user_id = u.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            JOIN kwv_workplaces w ON a.workplace_id = w.id
            {where} ORDER BY a.check_time DESC LIMIT %s OFFSET %s
        """, params + [per_page, offset])
        items = [{
            "id": r[0], "user_id": r[1], "user_name": r[2], "profile_photo": r[3],
            "nationality": r[4], "workplace_id": r[5], "workplace_name": r[6],
            "check_type": r[7], "check_method": r[8],
            "check_time": r[9].isoformat() if r[9] else None,
            "is_valid": bool(r[10]), "distance": r[11], "invalid_reason": r[12],
            "latitude": float(r[13]) if r[13] else None,
            "longitude": float(r[14]) if r[14] else None
        } for r in cursor.fetchall()]
        return {"items": items, "total": total, "page": page, "per_page": per_page}
    finally:
        conn.close()

@router.get("/attendance/admin/summary")
async def attendance_summary(
    date: Optional[str] = None,
    workplace_id: Optional[int] = None,
    user: dict = Depends(get_current_user)
):
    """출퇴근 요약 (일별)"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {}
    try:
        cursor = conn.cursor()
        target_date = date or datetime.utcnow().strftime('%Y-%m-%d')
        wp_filter = "AND a.workplace_id = %s" if workplace_id else ""
        params_base = [target_date] + ([workplace_id] if workplace_id else [])

        # 총 출근자 수
        cursor.execute(f"""
            SELECT COUNT(DISTINCT a.user_id)
            FROM kwv_attendance a
            WHERE DATE(a.check_time) = %s AND a.check_type = 'check_in' {wp_filter}
        """, params_base)
        total_checkins = cursor.fetchone()[0]

        # 유효/무효
        cursor.execute(f"""
            SELECT a.is_valid, COUNT(*)
            FROM kwv_attendance a
            WHERE DATE(a.check_time) = %s {wp_filter}
            GROUP BY a.is_valid
        """, params_base)
        valid_count, invalid_count = 0, 0
        for r in cursor.fetchall():
            if r[0]:
                valid_count = r[1]
            else:
                invalid_count = r[1]

        # 퇴근자 수
        cursor.execute(f"""
            SELECT COUNT(DISTINCT a.user_id)
            FROM kwv_attendance a
            WHERE DATE(a.check_time) = %s AND a.check_type = 'check_out' {wp_filter}
        """, params_base)
        total_checkouts = cursor.fetchone()[0]

        # 총 등록 근로자 (활성 배정)
        if workplace_id:
            cursor.execute("SELECT COUNT(*) FROM kwv_worker_assignments WHERE workplace_id = %s AND status = 'active'", (workplace_id,))
        else:
            cursor.execute("SELECT COUNT(*) FROM kwv_worker_assignments WHERE status = 'active'")
        total_workers = cursor.fetchone()[0]

        return {
            "date": target_date,
            "total_workers": total_workers,
            "total_checkins": total_checkins,
            "total_checkouts": total_checkouts,
            "valid_records": valid_count,
            "invalid_records": invalid_count,
            "attendance_rate": round(total_checkins / total_workers * 100, 1) if total_workers > 0 else 0
        }
    finally:
        conn.close()

@router.post("/attendance/admin/manual")
async def admin_manual_check(request: Request, user: dict = Depends(get_current_user)):
    """관리자 수동 출퇴근 등록"""
    require_admin_level(user, 2)
    body = await request.json()
    target_user_id = body.get("user_id")
    workplace_id = body.get("workplace_id")
    check_type = body.get("check_type", "check_in")
    note = body.get("note", "관리자 수동 등록")
    if not target_user_id or not workplace_id:
        raise HTTPException(status_code=400, detail="user_id, workplace_id 필요")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO kwv_attendance (user_id, workplace_id, check_type, check_method, is_valid, note)
            VALUES (%s, %s, %s, 'admin', 1, %s)
        """, (target_user_id, workplace_id, check_type, note))
        conn.commit()
        return {"id": cursor.lastrowid, "message": "수동 출퇴근 등록 완료"}
    finally:
        conn.close()

# ==================== Phase 5: 활동일지 + 포인트 ====================

class ActivityCreate(BaseModel):
    workplace_id: Optional[int] = None
    activity_date: str
    activity_type: str = "work"
    title: str
    content: Optional[str] = None
    hours: float = 0
    photo_url: Optional[str] = None
    photo_url_2: Optional[str] = None
    photo_url_3: Optional[str] = None

# --- 활동일지 ---

@router.post("/activities")
async def create_activity(data: ActivityCreate, user: dict = Depends(get_current_user)):
    """활동일지 작성"""
    user_id = user.get("user_id") or user.get("sub")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO kwv_activity_logs (user_id, workplace_id, activity_date, activity_type,
                title, content, hours, photo_url, photo_url_2, photo_url_3)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
        """, (user_id, data.workplace_id, data.activity_date, data.activity_type,
              data.title, data.content, data.hours, data.photo_url, data.photo_url_2, data.photo_url_3))
        activity_id = cursor.lastrowid
        # 포인트 자동 적립 (활동일지 작성)
        add_points(cursor, int(user_id), 'activity', 'activity_log', activity_id, 'activity_submitted')
        conn.commit()
        return {"id": activity_id, "message": "활동일지가 등록되었습니다"}
    finally:
        conn.close()

@router.get("/activities/my")
async def my_activities(
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    status: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """내 활동일지"""
    user_id = user.get("user_id") or user.get("sub")
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        where = "WHERE a.user_id = %s"
        params = [user_id]
        if date_from:
            where += " AND a.activity_date >= %s"
            params.append(date_from)
        if date_to:
            where += " AND a.activity_date <= %s"
            params.append(date_to)
        if status:
            where += " AND a.status = %s"
            params.append(status)
        cursor.execute(f"""
            SELECT a.id, a.activity_date, a.activity_type, a.title, a.content, a.hours,
                   a.photo_url, a.status, a.created_at, w.name as workplace_name
            FROM kwv_activity_logs a
            LEFT JOIN kwv_workplaces w ON a.workplace_id = w.id
            {where} ORDER BY a.activity_date DESC LIMIT 100
        """, params)
        return [{
            "id": r[0], "activity_date": r[1].isoformat() if r[1] else None,
            "activity_type": r[2], "title": r[3], "content": r[4], "hours": float(r[5]) if r[5] else 0,
            "photo_url": r[6], "status": r[7],
            "created_at": r[8].isoformat() if r[8] else None,
            "workplace_name": r[9]
        } for r in cursor.fetchall()]
    finally:
        conn.close()

@router.get("/activities/admin")
async def admin_activities(
    user_id: Optional[int] = None,
    workplace_id: Optional[int] = None,
    status: Optional[str] = None,
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    page: int = 1,
    per_page: int = 30,
    user: dict = Depends(get_current_user)
):
    """관리자 활동일지 조회"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": [], "total": 0}
    try:
        cursor = conn.cursor()
        where = "WHERE 1=1"
        params = []
        if user_id:
            where += " AND a.user_id = %s"
            params.append(user_id)
        if workplace_id:
            where += " AND a.workplace_id = %s"
            params.append(workplace_id)
        if status:
            where += " AND a.status = %s"
            params.append(status)
        if date_from:
            where += " AND a.activity_date >= %s"
            params.append(date_from)
        if date_to:
            where += " AND a.activity_date <= %s"
            params.append(date_to)

        cursor.execute(f"SELECT COUNT(*) FROM kwv_activity_logs a {where}", params)
        total = cursor.fetchone()[0]
        offset = (page - 1) * per_page
        cursor.execute(f"""
            SELECT a.id, a.user_id, u.name as user_name, u.profile_photo,
                   va.nationality, a.workplace_id, w.name as workplace_name,
                   a.activity_date, a.activity_type, a.title, a.content, a.hours,
                   a.photo_url, a.status, a.created_at
            FROM kwv_activity_logs a
            JOIN kwv_users u ON a.user_id = u.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            LEFT JOIN kwv_workplaces w ON a.workplace_id = w.id
            {where} ORDER BY a.activity_date DESC, a.created_at DESC LIMIT %s OFFSET %s
        """, params + [per_page, offset])
        items = [{
            "id": r[0], "user_id": r[1], "user_name": r[2], "profile_photo": r[3],
            "nationality": r[4], "workplace_id": r[5], "workplace_name": r[6],
            "activity_date": r[7].isoformat() if r[7] else None,
            "activity_type": r[8], "title": r[9], "content": r[10],
            "hours": float(r[11]) if r[11] else 0, "photo_url": r[12],
            "status": r[13], "created_at": r[14].isoformat() if r[14] else None
        } for r in cursor.fetchall()]
        return {"items": items, "total": total, "page": page, "per_page": per_page}
    finally:
        conn.close()

@router.put("/activities/{activity_id}/approve")
async def approve_activity(activity_id: int, user: dict = Depends(get_current_user)):
    """활동일지 승인"""
    require_admin_level(user, 2)
    admin_id = user.get("user_id") or user.get("sub")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            UPDATE kwv_activity_logs SET status = 'approved', approved_by = %s, approved_at = NOW()
            WHERE id = %s AND status = 'submitted'
        """, (admin_id, activity_id))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="활동일지를 찾을 수 없거나 이미 처리되었습니다")
        # 활동일지 작성자에게 승인 포인트
        cursor.execute("SELECT user_id FROM kwv_activity_logs WHERE id = %s", (activity_id,))
        row = cursor.fetchone()
        if row:
            add_points(cursor, row[0], 'activity', 'activity_log', activity_id, 'activity_approved')
        conn.commit()
        return {"message": "활동일지가 승인되었습니다"}
    finally:
        conn.close()

@router.put("/activities/{activity_id}/reject")
async def reject_activity(activity_id: int, request: Request, user: dict = Depends(get_current_user)):
    """활동일지 반려"""
    require_admin_level(user, 2)
    body = await request.json()
    reason = body.get("reason", "")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            UPDATE kwv_activity_logs SET status = 'rejected', rejection_reason = %s
            WHERE id = %s AND status = 'submitted'
        """, (reason, activity_id))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="활동일지를 찾을 수 없거나 이미 처리되었습니다")
        conn.commit()
        return {"message": "활동일지가 반려되었습니다"}
    finally:
        conn.close()

# --- 포인트 시스템 ---

def add_points(cursor, user_id: int, point_type: str, ref_type: str, ref_id: int, rule_key: str, created_by: int = None):
    """포인트 자동 적립 헬퍼"""
    cursor.execute("SELECT points FROM kwv_point_rules WHERE rule_key = %s AND is_active = 1", (rule_key,))
    row = cursor.fetchone()
    if not row or row[0] == 0:
        return 0
    pts = row[0]
    cursor.execute("""
        INSERT INTO kwv_points (user_id, points, point_type, reference_type, reference_id, description, created_by)
        VALUES (%s, %s, %s, %s, %s, %s, %s)
    """, (user_id, pts, point_type, ref_type, ref_id, rule_key, created_by))
    return pts

@router.get("/points/my")
async def my_points(user: dict = Depends(get_current_user)):
    """내 포인트 현황"""
    user_id = user.get("user_id") or user.get("sub")
    conn = get_kwv_db_connection()
    if not conn:
        return {"total": 0, "history": []}
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT COALESCE(SUM(points),0) FROM kwv_points WHERE user_id = %s", (user_id,))
        total = cursor.fetchone()[0]
        cursor.execute("""
            SELECT id, points, point_type, description, created_at
            FROM kwv_points WHERE user_id = %s ORDER BY created_at DESC LIMIT 50
        """, (user_id,))
        history = [{
            "id": r[0], "points": r[1], "point_type": r[2],
            "description": r[3], "created_at": r[4].isoformat() if r[4] else None
        } for r in cursor.fetchall()]
        return {"total": total, "history": history}
    finally:
        conn.close()

@router.get("/points/ranking")
async def points_ranking(
    limit: int = 20,
    user: dict = Depends(get_current_user)
):
    """포인트 랭킹"""
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute(f"""
            SELECT p.user_id, u.name, u.profile_photo, va.nationality,
                   SUM(p.points) as total_points
            FROM kwv_points p
            JOIN kwv_users u ON p.user_id = u.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            GROUP BY p.user_id ORDER BY total_points DESC LIMIT %s
        """, (limit,))
        rank = 1
        result = []
        for r in cursor.fetchall():
            result.append({
                "rank": rank, "user_id": r[0], "name": r[1], "profile_photo": r[2],
                "nationality": r[3], "total_points": int(r[4])
            })
            rank += 1
        return result
    finally:
        conn.close()

@router.get("/points/admin")
async def admin_points(
    user_id: Optional[int] = None,
    point_type: Optional[str] = None,
    page: int = 1,
    per_page: int = 50,
    user: dict = Depends(get_current_user)
):
    """관리자 포인트 조회"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": [], "total": 0}
    try:
        cursor = conn.cursor()
        where = "WHERE 1=1"
        params = []
        if user_id:
            where += " AND p.user_id = %s"
            params.append(user_id)
        if point_type:
            where += " AND p.point_type = %s"
            params.append(point_type)
        cursor.execute(f"SELECT COUNT(*) FROM kwv_points p {where}", params)
        total = cursor.fetchone()[0]
        offset = (page - 1) * per_page
        cursor.execute(f"""
            SELECT p.id, p.user_id, u.name as user_name, p.points, p.point_type,
                   p.description, p.created_at
            FROM kwv_points p
            JOIN kwv_users u ON p.user_id = u.id
            {where} ORDER BY p.created_at DESC LIMIT %s OFFSET %s
        """, params + [per_page, offset])
        items = [{
            "id": r[0], "user_id": r[1], "user_name": r[2], "points": r[3],
            "point_type": r[4], "description": r[5],
            "created_at": r[6].isoformat() if r[6] else None
        } for r in cursor.fetchall()]
        return {"items": items, "total": total, "page": page, "per_page": per_page}
    finally:
        conn.close()

@router.post("/points/admin/adjust")
async def admin_adjust_points(request: Request, user: dict = Depends(get_current_user)):
    """관리자 포인트 수동 조정"""
    require_admin_level(user, 2)
    body = await request.json()
    target_user_id = body.get("user_id")
    points = body.get("points", 0)
    description = body.get("description", "관리자 수동 조정")
    if not target_user_id or not points:
        raise HTTPException(status_code=400, detail="user_id, points 필요")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        admin_id = user.get("user_id") or user.get("sub")
        cursor.execute("""
            INSERT INTO kwv_points (user_id, points, point_type, description, created_by)
            VALUES (%s, %s, 'admin', %s, %s)
        """, (target_user_id, points, description, admin_id))
        conn.commit()
        return {"message": f"{points}포인트가 조정되었습니다"}
    finally:
        conn.close()

@router.get("/point-rules")
async def get_point_rules(user: dict = Depends(get_current_user)):
    """포인트 규칙 조회"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return []
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT id, rule_key, rule_name, points, description, is_active FROM kwv_point_rules ORDER BY id")
        return [{
            "id": r[0], "rule_key": r[1], "rule_name": r[2], "points": r[3],
            "description": r[4], "is_active": bool(r[5])
        } for r in cursor.fetchall()]
    finally:
        conn.close()

@router.put("/point-rules/{rule_id}")
async def update_point_rule(rule_id: int, request: Request, user: dict = Depends(get_current_user)):
    """포인트 규칙 수정"""
    require_admin_level(user, 9)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        allowed = ['points', 'rule_name', 'description', 'is_active']
        sets, params = [], []
        for k, v in body.items():
            if k in allowed:
                sets.append(f"{k} = %s")
                params.append(v)
        if not sets:
            raise HTTPException(status_code=400, detail="수정할 항목이 없습니다")
        params.append(rule_id)
        cursor.execute(f"UPDATE kwv_point_rules SET {', '.join(sets)} WHERE id = %s", params)
        conn.commit()
        return {"message": "포인트 규칙이 수정되었습니다"}
    finally:
        conn.close()

# ==================== Phase 6: 상담일지 ====================

class CounselingCreate(BaseModel):
    user_id: int
    counseling_date: str
    counseling_type: str = "in_person"
    category: str = "other"
    title: str
    content: Optional[str] = None
    action_taken: Optional[str] = None
    follow_up_date: Optional[str] = None
    follow_up_note: Optional[str] = None
    severity: str = "low"
    is_confidential: bool = False

@router.get("/counseling")
async def list_counseling(
    user_id: Optional[int] = None,
    category: Optional[str] = None,
    severity: Optional[str] = None,
    status: Optional[str] = None,
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    page: int = 1,
    per_page: int = 30,
    user: dict = Depends(get_current_user)
):
    """상담일지 목록"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": [], "total": 0}
    try:
        cursor = conn.cursor()
        where = "WHERE 1=1"
        params = []
        if user_id:
            where += " AND c.user_id = %s"
            params.append(user_id)
        if category:
            where += " AND c.category = %s"
            params.append(category)
        if severity:
            where += " AND c.severity = %s"
            params.append(severity)
        if status:
            where += " AND c.status = %s"
            params.append(status)
        if date_from:
            where += " AND c.counseling_date >= %s"
            params.append(date_from)
        if date_to:
            where += " AND c.counseling_date <= %s"
            params.append(date_to)
        cursor.execute(f"SELECT COUNT(*) FROM kwv_counseling c {where}", params)
        total = cursor.fetchone()[0]
        offset = (page - 1) * per_page
        cursor.execute(f"""
            SELECT c.id, c.user_id, u.name as user_name, va.nationality,
                   c.counselor_id, co.name as counselor_name,
                   c.counseling_date, c.counseling_type, c.category, c.title,
                   c.severity, c.status, c.follow_up_date, c.is_confidential, c.created_at
            FROM kwv_counseling c
            JOIN kwv_users u ON c.user_id = u.id
            JOIN kwv_users co ON c.counselor_id = co.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            {where} ORDER BY c.counseling_date DESC, c.created_at DESC LIMIT %s OFFSET %s
        """, params + [per_page, offset])
        items = [{
            "id": r[0], "user_id": r[1], "user_name": r[2], "nationality": r[3],
            "counselor_id": r[4], "counselor_name": r[5],
            "counseling_date": r[6].isoformat() if r[6] else None,
            "counseling_type": r[7], "category": r[8], "title": r[9],
            "severity": r[10], "status": r[11],
            "follow_up_date": r[12].isoformat() if r[12] else None,
            "is_confidential": bool(r[13]),
            "created_at": r[14].isoformat() if r[14] else None
        } for r in cursor.fetchall()]
        return {"items": items, "total": total, "page": page, "per_page": per_page}
    finally:
        conn.close()

@router.get("/counseling/summary")
async def counseling_summary(user: dict = Depends(get_current_user)):
    """상담 통계"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {}
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT COUNT(*) FROM kwv_counseling")
        total = cursor.fetchone()[0]
        cursor.execute("SELECT status, COUNT(*) FROM kwv_counseling GROUP BY status")
        by_status = {r[0]: r[1] for r in cursor.fetchall()}
        cursor.execute("SELECT category, COUNT(*) FROM kwv_counseling GROUP BY category ORDER BY COUNT(*) DESC")
        by_category = {r[0]: r[1] for r in cursor.fetchall()}
        cursor.execute("SELECT severity, COUNT(*) FROM kwv_counseling GROUP BY severity")
        by_severity = {r[0]: r[1] for r in cursor.fetchall()}
        cursor.execute("SELECT COUNT(*) FROM kwv_counseling WHERE follow_up_date IS NOT NULL AND follow_up_date <= CURDATE() AND status IN ('open','in_progress')")
        overdue = cursor.fetchone()[0]
        return {
            "total": total, "by_status": by_status, "by_category": by_category,
            "by_severity": by_severity, "overdue_followups": overdue
        }
    finally:
        conn.close()

@router.get("/counseling/{counsel_id}")
async def get_counseling(counsel_id: int, user: dict = Depends(get_current_user)):
    """상담일지 상세"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT c.*, u.name as user_name, co.name as counselor_name, va.nationality
            FROM kwv_counseling c
            JOIN kwv_users u ON c.user_id = u.id
            JOIN kwv_users co ON c.counselor_id = co.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            WHERE c.id = %s
        """, (counsel_id,))
        r = cursor.fetchone()
        if not r:
            raise HTTPException(status_code=404, detail="상담일지를 찾을 수 없습니다")
        cols = [d[0] for d in cursor.description]
        result = {}
        for i, col in enumerate(cols):
            val = r[i]
            if hasattr(val, 'isoformat'):
                val = val.isoformat()
            if isinstance(val, bytes):
                val = bool(val[0]) if len(val) == 1 else val.decode()
            result[col] = val
        return result
    finally:
        conn.close()

@router.post("/counseling")
async def create_counseling(data: CounselingCreate, user: dict = Depends(get_current_user)):
    """상담일지 등록"""
    require_admin_level(user, 2)
    counselor_id = user.get("user_id") or user.get("sub")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO kwv_counseling (user_id, counselor_id, counseling_date, counseling_type,
                category, title, content, action_taken, follow_up_date, follow_up_note,
                severity, is_confidential)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
        """, (data.user_id, counselor_id, data.counseling_date, data.counseling_type,
              data.category, data.title, data.content, data.action_taken,
              data.follow_up_date or None, data.follow_up_note,
              data.severity, data.is_confidential))
        conn.commit()
        return {"id": cursor.lastrowid, "message": "상담일지가 등록되었습니다"}
    finally:
        conn.close()

@router.put("/counseling/{counsel_id}")
async def update_counseling(counsel_id: int, request: Request, user: dict = Depends(get_current_user)):
    """상담일지 수정"""
    require_admin_level(user, 2)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        allowed = ['counseling_date','counseling_type','category','title','content',
                    'action_taken','follow_up_date','follow_up_note','severity','status','is_confidential']
        sets, params = [], []
        for k, v in body.items():
            if k in allowed:
                sets.append(f"{k} = %s")
                params.append(v)
        if not sets:
            raise HTTPException(status_code=400, detail="수정할 항목이 없습니다")
        params.append(counsel_id)
        cursor.execute(f"UPDATE kwv_counseling SET {', '.join(sets)} WHERE id = %s", params)
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="상담일지를 찾을 수 없습니다")
        conn.commit()
        return {"message": "상담일지가 수정되었습니다"}
    finally:
        conn.close()

@router.delete("/counseling/{counsel_id}")
async def delete_counseling(counsel_id: int, user: dict = Depends(get_current_user)):
    """상담일지 삭제"""
    require_admin_level(user, 9)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM kwv_counseling WHERE id = %s", (counsel_id,))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="상담일지를 찾을 수 없습니다")
        conn.commit()
        return {"message": "상담일지가 삭제되었습니다"}
    finally:
        conn.close()

# ==================== Phase 7: 리포트 + Excel/PDF ====================

from fastapi.responses import StreamingResponse
import csv
import io

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

@router.get("/reports/attendance/excel")
async def export_attendance_excel(
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    workplace_id: Optional[int] = None,
    user: dict = Depends(get_current_user)
):
    """출퇴근 기록 Excel 다운로드"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        where = "WHERE 1=1"
        params = []
        if date_from:
            where += " AND DATE(a.check_time) >= %s"
            params.append(date_from)
        if date_to:
            where += " AND DATE(a.check_time) <= %s"
            params.append(date_to)
        if workplace_id:
            where += " AND a.workplace_id = %s"
            params.append(workplace_id)
        cursor.execute(f"""
            SELECT u.name, u.email, va.nationality, va.visa_type,
                   w.name as workplace, a.check_type, a.check_method,
                   a.check_time, a.is_valid, a.distance_from_workplace, a.invalid_reason
            FROM kwv_attendance a
            JOIN kwv_users u ON a.user_id = u.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            JOIN kwv_workplaces w ON a.workplace_id = w.id
            {where} ORDER BY a.check_time DESC
        """, params)
        rows = cursor.fetchall()

        if OPENPYXL_AVAILABLE:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "출퇴근기록"
            headers = ['이름','이메일','국적','비자','사업장','유형','방법','시간','유효','거리(m)','비고']
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True, size=11)
            thin_border = Border(
                left=Side(style='thin'), right=Side(style='thin'),
                top=Side(style='thin'), bottom=Side(style='thin')
            )
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border
            type_map = {'check_in':'출근','check_out':'퇴근'}
            method_map = {'qr':'QR','gps':'GPS','manual':'수동','admin':'관리자'}
            for i, r in enumerate(rows, 2):
                vals = [r[0], r[1], r[2], r[3], r[4],
                        type_map.get(r[5], r[5]), method_map.get(r[6], r[6]),
                        r[7].strftime('%Y-%m-%d %H:%M') if r[7] else '',
                        '유효' if r[8] else '이탈', r[9] or '', r[10] or '']
                for col, v in enumerate(vals, 1):
                    cell = ws.cell(row=i, column=col, value=v)
                    cell.border = thin_border
            for col in range(1, len(headers)+1):
                ws.column_dimensions[openpyxl.utils.get_column_letter(col)].width = 15
            output = io.BytesIO()
            wb.save(output)
            output.seek(0)
            filename = f"attendance_{date_from or 'all'}_{date_to or 'all'}.xlsx"
            return StreamingResponse(output,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                headers={"Content-Disposition": f"attachment; filename={filename}"})
        else:
            # CSV fallback
            output = io.StringIO()
            writer = csv.writer(output)
            writer.writerow(['이름','이메일','국적','비자','사업장','유형','방법','시간','유효','거리','비고'])
            for r in rows:
                writer.writerow([r[0], r[1], r[2], r[3], r[4], r[5], r[6],
                    r[7].strftime('%Y-%m-%d %H:%M') if r[7] else '', '유효' if r[8] else '이탈', r[9] or '', r[10] or ''])
            return StreamingResponse(io.BytesIO(output.getvalue().encode('utf-8-sig')),
                media_type="text/csv",
                headers={"Content-Disposition": "attachment; filename=attendance.csv"})
    finally:
        conn.close()

@router.get("/reports/applicants/excel")
async def export_applicants_excel(
    is_approved: Optional[str] = None,
    nationality: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """지원자 목록 Excel 다운로드"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        where = "WHERE u.user_type = 'applicant'"
        params = []
        if is_approved == 'true':
            where += " AND u.is_approved = 1"
        elif is_approved == 'false':
            where += " AND u.is_approved = 0"
        if nationality:
            where += " AND va.nationality = %s"
            params.append(nationality)
        cursor.execute(f"""
            SELECT u.id, u.name, u.email, u.phone, va.nationality, va.visa_type,
                   va.passport_number, va.birth_date, va.gender,
                   u.is_approved, lg.name as lg_name, u.created_at
            FROM kwv_users u
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            LEFT JOIN kwv_local_governments lg ON u.local_government_id = lg.id
            {where} ORDER BY u.created_at DESC
        """, params)
        rows = cursor.fetchall()

        if OPENPYXL_AVAILABLE:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "지원자목록"
            headers = ['ID','이름','이메일','전화번호','국적','비자','여권번호','생년월일','성별','승인','배정지자체','가입일']
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True, size=11)
            thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border
            gender_map = {'male':'남','female':'여'}
            for i, r in enumerate(rows, 2):
                vals = [r[0], r[1], r[2], r[3], r[4], r[5], r[6],
                        r[7].strftime('%Y-%m-%d') if r[7] else '',
                        gender_map.get(r[8], r[8] or ''),
                        '승인' if r[9] else '대기', r[10] or '',
                        r[11].strftime('%Y-%m-%d') if r[11] else '']
                for col, v in enumerate(vals, 1):
                    cell = ws.cell(row=i, column=col, value=v)
                    cell.border = thin_border
            for col in range(1, len(headers)+1):
                ws.column_dimensions[openpyxl.utils.get_column_letter(col)].width = 15
            output = io.BytesIO()
            wb.save(output)
            output.seek(0)
            return StreamingResponse(output,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                headers={"Content-Disposition": "attachment; filename=applicants.xlsx"})
        else:
            output = io.StringIO()
            writer = csv.writer(output)
            writer.writerow(headers)
            for r in rows:
                writer.writerow(r)
            return StreamingResponse(io.BytesIO(output.getvalue().encode('utf-8-sig')),
                media_type="text/csv",
                headers={"Content-Disposition": "attachment; filename=applicants.csv"})
    finally:
        conn.close()

@router.get("/reports/points/excel")
async def export_points_excel(user: dict = Depends(get_current_user)):
    """포인트 내역 Excel 다운로드"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT u.name, u.email, va.nationality, p.points, p.point_type,
                   p.description, p.created_at
            FROM kwv_points p
            JOIN kwv_users u ON p.user_id = u.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            ORDER BY p.created_at DESC
        """)
        rows = cursor.fetchall()
        if OPENPYXL_AVAILABLE:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "포인트내역"
            headers = ['이름','이메일','국적','포인트','유형','설명','일시']
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True, size=11)
            thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.fill = header_fill; cell.font = header_font
                cell.alignment = Alignment(horizontal='center'); cell.border = thin_border
            type_map = {'attendance':'출퇴근','activity':'활동','training':'교육','community':'봉사','bonus':'보너스','penalty':'패널티','admin':'관리자'}
            for i, r in enumerate(rows, 2):
                vals = [r[0], r[1], r[2], r[3], type_map.get(r[4],r[4]), r[5],
                        r[6].strftime('%Y-%m-%d %H:%M') if r[6] else '']
                for col, v in enumerate(vals, 1):
                    cell = ws.cell(row=i, column=col, value=v); cell.border = thin_border
            for col in range(1, len(headers)+1):
                ws.column_dimensions[openpyxl.utils.get_column_letter(col)].width = 15
            output = io.BytesIO()
            wb.save(output); output.seek(0)
            return StreamingResponse(output,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                headers={"Content-Disposition": "attachment; filename=points.xlsx"})
        else:
            output = io.StringIO()
            writer = csv.writer(output)
            writer.writerow(['이름','이메일','국적','포인트','유형','설명','일시'])
            for r in rows: writer.writerow(r)
            return StreamingResponse(io.BytesIO(output.getvalue().encode('utf-8-sig')),
                media_type="text/csv", headers={"Content-Disposition": "attachment; filename=points.csv"})
    finally:
        conn.close()

@router.get("/reports/counseling/excel")
async def export_counseling_excel(user: dict = Depends(get_current_user)):
    """상담일지 Excel 다운로드"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=503, detail="DB connection failed")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT u.name, va.nationality, co.name as counselor,
                   c.counseling_date, c.counseling_type, c.category,
                   c.title, c.content, c.action_taken, c.severity, c.status,
                   c.follow_up_date, c.is_confidential
            FROM kwv_counseling c
            JOIN kwv_users u ON c.user_id = u.id
            JOIN kwv_users co ON c.counselor_id = co.id
            LEFT JOIN kwv_visa_applicants va ON va.user_id = u.id
            ORDER BY c.counseling_date DESC
        """)
        rows = cursor.fetchall()
        if OPENPYXL_AVAILABLE:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "상담일지"
            headers = ['근로자','국적','상담사','날짜','방식','분류','제목','내용','조치','심각도','상태','후속상담일','기밀']
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True, size=11)
            thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.fill = header_fill; cell.font = header_font
                cell.alignment = Alignment(horizontal='center'); cell.border = thin_border
            cat_map = {'work':'근무','health':'건강','legal':'법률','housing':'숙소','salary':'급여','homesick':'향수병','conflict':'갈등','other':'기타'}
            sev_map = {'low':'낮음','medium':'보통','high':'높음','urgent':'긴급'}
            st_map = {'open':'진행중','in_progress':'처리중','resolved':'해결','closed':'종료'}
            type_map = {'in_person':'대면','phone':'전화','video':'화상','text':'문자'}
            for i, r in enumerate(rows, 2):
                vals = [r[0], r[1], r[2],
                        r[3].strftime('%Y-%m-%d') if r[3] else '',
                        type_map.get(r[4],r[4]), cat_map.get(r[5],r[5]),
                        r[6], r[7], r[8], sev_map.get(r[9],r[9]), st_map.get(r[10],r[10]),
                        r[11].strftime('%Y-%m-%d') if r[11] else '',
                        'Y' if r[12] else 'N']
                for col, v in enumerate(vals, 1):
                    cell = ws.cell(row=i, column=col, value=v); cell.border = thin_border
            for col in range(1, len(headers)+1):
                ws.column_dimensions[openpyxl.utils.get_column_letter(col)].width = 16
            output = io.BytesIO()
            wb.save(output); output.seek(0)
            return StreamingResponse(output,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                headers={"Content-Disposition": "attachment; filename=counseling.xlsx"})
        else:
            output = io.StringIO()
            writer = csv.writer(output)
            writer.writerow(headers)
            for r in rows: writer.writerow(r)
            return StreamingResponse(io.BytesIO(output.getvalue().encode('utf-8-sig')),
                media_type="text/csv", headers={"Content-Disposition": "attachment; filename=counseling.csv"})
    finally:
        conn.close()

@router.get("/reports/dashboard")
async def dashboard_report(user: dict = Depends(get_current_user)):
    """종합 대시보드 리포트 데이터"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {}
    try:
        cursor = conn.cursor()
        # 총 근로자
        cursor.execute("SELECT COUNT(*) FROM kwv_users WHERE user_type='applicant'")
        total_workers = cursor.fetchone()[0]
        cursor.execute("SELECT COUNT(*) FROM kwv_users WHERE user_type='applicant' AND is_approved=1")
        approved_workers = cursor.fetchone()[0]
        # 총 사업장
        cursor.execute("SELECT COUNT(*) FROM kwv_workplaces WHERE is_active=1")
        total_workplaces = cursor.fetchone()[0]
        # 총 지자체
        cursor.execute("SELECT COUNT(*) FROM kwv_local_governments WHERE is_active=1")
        total_lgs = cursor.fetchone()[0]
        # MOU
        cursor.execute("SELECT COUNT(*) FROM kwv_mou_agreements WHERE status='active'")
        active_mous = cursor.fetchone()[0]
        # 오늘 출근
        cursor.execute("SELECT COUNT(DISTINCT user_id) FROM kwv_attendance WHERE check_type='check_in' AND DATE(check_time)=CURDATE()")
        today_checkins = cursor.fetchone()[0]
        # 포인트 총합
        cursor.execute("SELECT COALESCE(SUM(points),0) FROM kwv_points")
        total_points = cursor.fetchone()[0]
        # 상담
        cursor.execute("SELECT COUNT(*) FROM kwv_counseling WHERE status IN ('open','in_progress')")
        open_counseling = cursor.fetchone()[0]
        return {
            "total_workers": total_workers,
            "approved_workers": approved_workers,
            "total_workplaces": total_workplaces,
            "total_lgs": total_lgs,
            "active_mous": active_mous,
            "today_checkins": today_checkins,
            "total_points_issued": int(total_points),
            "open_counseling": open_counseling
        }
    finally:
        conn.close()

# ==================== Phase 8: 이상감지 + 알림 ====================

# --- 이상감지 스캔 ---
@router.post("/admin/anomalies/scan")
async def scan_anomalies(user: dict = Depends(get_current_user)):
    """이상감지 규칙 기반 스캔"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        # 활성 규칙 로드
        cursor.execute("SELECT * FROM kwv_anomaly_rules WHERE is_active=1")
        rules = {r['rule_key']: r for r in cursor.fetchall()}

        detected = []
        today = datetime.now().date()

        # 규칙1: 3일 연속 결근
        rule = rules.get('absent_3days')
        if rule:
            threshold = rule['threshold'] or 3
            cursor.execute("""
                SELECT u.id, u.name, u.email,
                    (SELECT MAX(DATE(check_time)) FROM kwv_attendance WHERE user_id=u.id AND check_type='check_in') as last_checkin
                FROM kwv_users u
                WHERE u.user_type='applicant' AND u.is_active=1 AND u.is_approved=1
            """)
            workers = cursor.fetchall()
            for w in workers:
                last = w['last_checkin']
                if last and (today - last).days >= threshold:
                    days_absent = (today - last).days
                    # 중복 체크
                    cursor.execute("""
                        SELECT id FROM kwv_anomalies WHERE user_id=%s AND anomaly_type='absent_streak'
                        AND status IN ('detected','reviewing') AND DATE(created_at)=CURDATE()
                    """, (w['id'],))
                    if not cursor.fetchone():
                        cursor.execute("""
                            INSERT INTO kwv_anomalies (user_id, anomaly_type, score, description, details)
                            VALUES (%s, 'absent_streak', %s, %s, %s)
                        """, (w['id'], rule['score'],
                              f"{w['name']} - {days_absent}일 연속 결근",
                              json.dumps({"days_absent": days_absent, "last_checkin": str(last)})))
                        detected.append({"user": w['name'], "type": "absent_streak", "days": days_absent})

        # 규칙2: GPS 이탈 (최근 7일)
        rule = rules.get('gps_violation')
        if rule:
            cursor.execute("""
                SELECT a.user_id, u.name, COUNT(*) as cnt
                FROM kwv_attendance a
                JOIN kwv_users u ON u.id=a.user_id
                WHERE a.is_valid=0 AND a.check_time >= DATE_SUB(NOW(), INTERVAL 7 DAY)
                GROUP BY a.user_id
                HAVING cnt >= 2
            """)
            for row in cursor.fetchall():
                cursor.execute("""
                    SELECT id FROM kwv_anomalies WHERE user_id=%s AND anomaly_type='gps_violation'
                    AND status IN ('detected','reviewing') AND created_at >= DATE_SUB(NOW(), INTERVAL 7 DAY)
                """, (row['user_id'],))
                if not cursor.fetchone():
                    cursor.execute("""
                        INSERT INTO kwv_anomalies (user_id, anomaly_type, score, description, details)
                        VALUES (%s, 'gps_violation', %s, %s, %s)
                    """, (row['user_id'], rule['score'],
                          f"{row['name']} - 최근 7일 GPS 이탈 {row['cnt']}회",
                          json.dumps({"violation_count": row['cnt']})))
                    detected.append({"user": row['name'], "type": "gps_violation", "count": row['cnt']})

        # 규칙3: 미퇴근 3회 이상 (최근 14일)
        rule = rules.get('no_checkout_3')
        if rule:
            threshold = rule['threshold'] or 3
            cursor.execute("""
                SELECT ci.user_id, u.name, COUNT(*) as cnt
                FROM kwv_attendance ci
                JOIN kwv_users u ON u.id=ci.user_id
                LEFT JOIN kwv_attendance co ON co.user_id=ci.user_id
                    AND co.check_type='check_out' AND DATE(co.check_time)=DATE(ci.check_time)
                WHERE ci.check_type='check_in' AND ci.check_time >= DATE_SUB(NOW(), INTERVAL 14 DAY)
                    AND co.id IS NULL
                GROUP BY ci.user_id
                HAVING cnt >= %s
            """, (threshold,))
            for row in cursor.fetchall():
                cursor.execute("""
                    SELECT id FROM kwv_anomalies WHERE user_id=%s AND anomaly_type='no_checkout'
                    AND status IN ('detected','reviewing') AND created_at >= DATE_SUB(NOW(), INTERVAL 14 DAY)
                """, (row['user_id'],))
                if not cursor.fetchone():
                    cursor.execute("""
                        INSERT INTO kwv_anomalies (user_id, anomaly_type, score, description, details)
                        VALUES (%s, 'no_checkout', %s, %s, %s)
                    """, (row['user_id'], rule['score'],
                          f"{row['name']} - 최근 14일 미퇴근 {row['cnt']}회",
                          json.dumps({"no_checkout_count": row['cnt']})))
                    detected.append({"user": row['name'], "type": "no_checkout", "count": row['cnt']})

        conn.commit()

        # 감지된 이상에 대해 관리자 알림 생성
        if detected:
            cursor.execute("SELECT id FROM kwv_users WHERE admin_level >= 5")
            admins = cursor.fetchall()
            for admin in admins:
                cursor.execute("""
                    INSERT INTO kwv_notifications (user_id, title, message, notification_type, reference_type)
                    VALUES (%s, %s, %s, 'warning', 'anomaly')
                """, (admin['id'],
                      f"이상감지 스캔: {len(detected)}건 발견",
                      f"스캔 결과 {len(detected)}건의 이상 행동이 감지되었습니다."))
            conn.commit()

        return {"scanned": True, "detected_count": len(detected), "details": detected}
    finally:
        conn.close()

@router.get("/admin/anomalies")
async def list_anomalies(
    status: str = None, anomaly_type: str = None,
    page: int = 1, per_page: int = 20,
    user: dict = Depends(get_current_user)
):
    """이상감지 목록"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": [], "total": 0}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        where = ["1=1"]
        params = []
        if status:
            where.append("a.status=%s")
            params.append(status)
        if anomaly_type:
            where.append("a.anomaly_type=%s")
            params.append(anomaly_type)

        w = " AND ".join(where)
        cursor.execute(f"SELECT COUNT(*) as cnt FROM kwv_anomalies a WHERE {w}", params)
        total = cursor.fetchone()['cnt']

        offset = (page - 1) * per_page
        cursor.execute(f"""
            SELECT a.*, u.name as user_name, u.email as user_email,
                   r.name as resolver_name
            FROM kwv_anomalies a
            JOIN kwv_users u ON u.id=a.user_id
            LEFT JOIN kwv_users r ON r.id=a.resolved_by
            WHERE {w}
            ORDER BY a.score DESC, a.created_at DESC
            LIMIT %s OFFSET %s
        """, params + [per_page, offset])
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
                if k == 'details' and isinstance(v, str):
                    try:
                        item[k] = json.loads(v)
                    except:
                        pass
        return {"items": items, "total": total, "page": page, "per_page": per_page}
    finally:
        conn.close()

@router.put("/admin/anomalies/{anomaly_id}/resolve")
async def resolve_anomaly(anomaly_id: int, body: dict = Body(...), user: dict = Depends(get_current_user)):
    """이상감지 해결 처리"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor()
        new_status = body.get("status", "resolved")
        note = body.get("note", "")
        cursor.execute("""
            UPDATE kwv_anomalies SET status=%s, resolved_by=%s, resolved_at=NOW(), resolve_note=%s
            WHERE id=%s
        """, (new_status, int(user['sub']), note, anomaly_id))
        conn.commit()
        return {"success": True}
    finally:
        conn.close()

@router.get("/admin/anomalies/summary")
async def anomaly_summary(user: dict = Depends(get_current_user)):
    """이상감지 요약"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("SELECT COUNT(*) as cnt FROM kwv_anomalies WHERE status IN ('detected','reviewing')")
        active = cursor.fetchone()['cnt']
        cursor.execute("SELECT COUNT(*) as cnt FROM kwv_anomalies WHERE score >= 30 AND status IN ('detected','reviewing')")
        high_risk = cursor.fetchone()['cnt']
        cursor.execute("""
            SELECT anomaly_type, COUNT(*) as cnt
            FROM kwv_anomalies WHERE status IN ('detected','reviewing')
            GROUP BY anomaly_type
        """)
        by_type = {r['anomaly_type']: r['cnt'] for r in cursor.fetchall()}
        cursor.execute("SELECT COUNT(*) as cnt FROM kwv_anomalies WHERE status='resolved'")
        resolved = cursor.fetchone()['cnt']
        return {"active": active, "high_risk": high_risk, "resolved": resolved, "by_type": by_type}
    finally:
        conn.close()

# --- 알림 ---
@router.get("/notifications/my")
async def my_notifications(page: int = 1, per_page: int = 20, user: dict = Depends(get_current_user)):
    """내 알림 목록"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": [], "total": 0}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("SELECT COUNT(*) as cnt FROM kwv_notifications WHERE user_id=%s", (int(user['sub']),))
        total = cursor.fetchone()['cnt']
        offset = (page - 1) * per_page
        cursor.execute("""
            SELECT * FROM kwv_notifications WHERE user_id=%s
            ORDER BY created_at DESC LIMIT %s OFFSET %s
        """, (int(user['sub']), per_page, offset))
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
        return {"items": items, "total": total}
    finally:
        conn.close()

@router.get("/notifications/my/unread-count")
async def unread_count(user: dict = Depends(get_current_user)):
    """안읽은 알림 수"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"count": 0}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("SELECT COUNT(*) as cnt FROM kwv_notifications WHERE user_id=%s AND is_read=0", (int(user['sub']),))
        return {"count": cursor.fetchone()['cnt']}
    finally:
        conn.close()

@router.put("/notifications/{noti_id}/read")
async def mark_notification_read(noti_id: int, user: dict = Depends(get_current_user)):
    """알림 읽음 처리"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor()
        cursor.execute("UPDATE kwv_notifications SET is_read=1, read_at=NOW() WHERE id=%s AND user_id=%s", (noti_id, int(user['sub'])))
        conn.commit()
        return {"success": True}
    finally:
        conn.close()

@router.put("/notifications/read-all")
async def mark_all_read(user: dict = Depends(get_current_user)):
    """모든 알림 읽음 처리"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor()
        cursor.execute("UPDATE kwv_notifications SET is_read=1, read_at=NOW() WHERE user_id=%s AND is_read=0", (int(user['sub']),))
        conn.commit()
        return {"success": True, "updated": cursor.rowcount}
    finally:
        conn.close()

@router.post("/admin/notifications/send")
async def send_notification(body: dict = Body(...), user: dict = Depends(get_current_user)):
    """관리자 알림 전송"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor()
        target = body.get("target", "all")  # all, admins, user_id
        title = body.get("title", "")
        message = body.get("message", "")
        ntype = body.get("type", "info")

        if not title:
            raise HTTPException(400, "제목을 입력하세요")

        user_ids = []
        if target == "all":
            cursor.execute("SELECT id FROM kwv_users WHERE is_active=1")
            user_ids = [r[0] for r in cursor.fetchall()]
        elif target == "admins":
            cursor.execute("SELECT id FROM kwv_users WHERE admin_level >= 5")
            user_ids = [r[0] for r in cursor.fetchall()]
        elif target == "workers":
            cursor.execute("SELECT id FROM kwv_users WHERE user_type='applicant' AND is_active=1")
            user_ids = [r[0] for r in cursor.fetchall()]
        else:
            try:
                user_ids = [int(target)]
            except:
                raise HTTPException(400, "유효하지 않은 대상")

        for uid in user_ids:
            cursor.execute("""
                INSERT INTO kwv_notifications (user_id, title, message, notification_type)
                VALUES (%s, %s, %s, %s)
            """, (uid, title, message, ntype))
        conn.commit()
        return {"success": True, "sent_count": len(user_ids)}
    finally:
        conn.close()

# ==================== Phase 9: 보험 관리 ====================

class InsuranceCreate(BaseModel):
    user_id: int
    insurance_type: str = "health"
    provider: Optional[str] = None
    policy_number: Optional[str] = None
    start_date: str
    end_date: str
    premium: Optional[int] = 0
    coverage: Optional[str] = None
    status: Optional[str] = "active"
    note: Optional[str] = None

@router.post("/admin/insurance")
async def create_insurance(ins: InsuranceCreate, user: dict = Depends(get_current_user)):
    """보험 등록"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO kwv_insurance (user_id, insurance_type, provider, policy_number,
                start_date, end_date, premium, coverage, status, note, created_by)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
        """, (ins.user_id, ins.insurance_type, ins.provider, ins.policy_number,
              ins.start_date, ins.end_date, ins.premium, ins.coverage, ins.status,
              ins.note, int(user['sub'])))
        conn.commit()
        return {"success": True, "id": cursor.lastrowid}
    finally:
        conn.close()

@router.put("/admin/insurance/{ins_id}")
async def update_insurance(ins_id: int, body: dict = Body(...), user: dict = Depends(get_current_user)):
    """보험 수정"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor()
        fields = []
        params = []
        for key in ['insurance_type','provider','policy_number','start_date','end_date','premium','coverage','status','note']:
            if key in body:
                fields.append(f"{key}=%s")
                params.append(body[key])
        if not fields:
            raise HTTPException(400, "수정할 항목이 없습니다")
        params.append(ins_id)
        cursor.execute(f"UPDATE kwv_insurance SET {','.join(fields)} WHERE id=%s", params)
        conn.commit()
        return {"success": True}
    finally:
        conn.close()

@router.delete("/admin/insurance/{ins_id}")
async def delete_insurance(ins_id: int, user: dict = Depends(get_current_user)):
    """보험 삭제"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(500, "DB 연결 실패")
    try:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM kwv_insurance WHERE id=%s", (ins_id,))
        conn.commit()
        return {"success": True}
    finally:
        conn.close()

@router.get("/admin/insurance")
async def list_insurance(
    user_id: int = None, insurance_type: str = None, status: str = None,
    page: int = 1, per_page: int = 20,
    user: dict = Depends(get_current_user)
):
    """보험 목록 (관리자)"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": [], "total": 0}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        where = ["1=1"]
        params = []
        if user_id:
            where.append("i.user_id=%s"); params.append(user_id)
        if insurance_type:
            where.append("i.insurance_type=%s"); params.append(insurance_type)
        if status:
            where.append("i.status=%s"); params.append(status)

        w = " AND ".join(where)
        cursor.execute(f"SELECT COUNT(*) as cnt FROM kwv_insurance i WHERE {w}", params)
        total = cursor.fetchone()['cnt']

        offset = (page - 1) * per_page
        cursor.execute(f"""
            SELECT i.*, u.name as user_name, u.email as user_email
            FROM kwv_insurance i
            JOIN kwv_users u ON u.id=i.user_id
            WHERE {w}
            ORDER BY i.end_date ASC
            LIMIT %s OFFSET %s
        """, params + [per_page, offset])
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
                if isinstance(v, type(None)):
                    pass
        return {"items": items, "total": total, "page": page}
    finally:
        conn.close()

@router.get("/admin/insurance/expiring")
async def expiring_insurance(days: int = 30, user: dict = Depends(get_current_user)):
    """만료 임박 보험"""
    require_admin(user)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": []}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            SELECT i.*, u.name as user_name, u.email as user_email
            FROM kwv_insurance i
            JOIN kwv_users u ON u.id=i.user_id
            WHERE i.status='active' AND i.end_date BETWEEN CURDATE() AND DATE_ADD(CURDATE(), INTERVAL %s DAY)
            ORDER BY i.end_date ASC
        """, (days,))
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
        return {"items": items, "total": len(items)}
    finally:
        conn.close()

@router.get("/insurance/my")
async def my_insurance(user: dict = Depends(get_current_user)):
    """내 보험 목록"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": []}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            SELECT * FROM kwv_insurance WHERE user_id=%s ORDER BY end_date DESC
        """, (int(user['sub']),))
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
        return {"items": items}
    finally:
        conn.close()

# ==================== 공지사항 (Notices) ====================

@router.get("/notices")
async def list_notices(
    target_type: Optional[str] = None,
    local_government_id: Optional[int] = None,
    important_only: Optional[int] = 0
):
    """공지사항 목록 (공개 - 로그인 불필요)"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": []}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        sql = """
            SELECT n.*, u.name as created_by_name,
                   lg.name as lg_name
            FROM kwv_notices n
            LEFT JOIN kwv_users u ON n.created_by = u.id
            LEFT JOIN kwv_local_governments lg ON n.local_government_id = lg.id
            WHERE n.is_active = 1
        """
        params = []
        if target_type:
            sql += " AND n.target_type = %s"
            params.append(target_type)
        if local_government_id:
            sql += " AND (n.target_type = 'all' OR n.local_government_id = %s)"
            params.append(local_government_id)
        if important_only:
            sql += " AND n.is_important = 1"
        sql += " ORDER BY n.is_important DESC, n.created_at DESC"
        cursor.execute(sql, params)
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
        return {"items": items, "total": len(items)}
    finally:
        conn.close()

@router.get("/notices/{notice_id}")
async def get_notice(notice_id: int):
    """공지사항 상세"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        # 조회수 증가
        cursor.execute("UPDATE kwv_notices SET view_count = view_count + 1 WHERE id = %s", (notice_id,))
        conn.commit()
        cursor.execute("""
            SELECT n.*, u.name as created_by_name, lg.name as lg_name
            FROM kwv_notices n
            LEFT JOIN kwv_users u ON n.created_by = u.id
            LEFT JOIN kwv_local_governments lg ON n.local_government_id = lg.id
            WHERE n.id = %s
        """, (notice_id,))
        item = cursor.fetchone()
        if not item:
            raise HTTPException(status_code=404, detail="공지를 찾을 수 없습니다")
        for k, v in item.items():
            if isinstance(v, (datetime, date)):
                item[k] = v.isoformat()
        return item
    finally:
        conn.close()

@router.post("/admin/notices")
async def create_notice(request: Request, user: dict = Depends(get_current_user)):
    """공지사항 작성 (관리자)"""
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            INSERT INTO kwv_notices (title, content, target_type, local_government_id, is_important, created_by)
            VALUES (%s, %s, %s, %s, %s, %s)
        """, (
            body['title'], body['content'],
            body.get('target_type', 'all'),
            body.get('local_government_id') or None,
            body.get('is_important', 0),
            int(user['sub'])
        ))
        conn.commit()
        return {"message": "공지가 등록되었습니다", "id": cursor.lastrowid}
    finally:
        conn.close()

@router.put("/admin/notices/{notice_id}")
async def update_notice(notice_id: int, request: Request, user: dict = Depends(get_current_user)):
    """공지사항 수정 (관리자)"""
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            UPDATE kwv_notices SET title=%s, content=%s, target_type=%s,
            local_government_id=%s, is_important=%s WHERE id=%s
        """, (
            body['title'], body['content'],
            body.get('target_type', 'all'),
            body.get('local_government_id') or None,
            body.get('is_important', 0),
            notice_id
        ))
        conn.commit()
        return {"message": "공지가 수정되었습니다"}
    finally:
        conn.close()

@router.delete("/admin/notices/{notice_id}")
async def delete_notice(notice_id: int, user: dict = Depends(get_current_user)):
    """공지사항 삭제 (관리자)"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("UPDATE kwv_notices SET is_active = 0 WHERE id = %s", (notice_id,))
        conn.commit()
        return {"message": "공지가 삭제되었습니다"}
    finally:
        conn.close()

# ==================== 구인 관리 API (Jobs) ====================

@router.get("/jobs")
async def list_jobs(
    visa_type: Optional[str] = None,
    local_government_id: Optional[int] = None,
    status: Optional[str] = None
):
    """구인 목록 (공개)"""
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": []}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        sql = """
            SELECT j.*, u.name as created_by_name,
                   lg.name as lg_name
            FROM kwv_jobs j
            LEFT JOIN kwv_users u ON j.created_by = u.id
            LEFT JOIN kwv_local_governments lg ON j.local_government_id = lg.id
            WHERE 1=1
        """
        params = []
        if status:
            sql += " AND j.status = %s"
            params.append(status)
        else:
            sql += " AND j.status = 'active'"
        if visa_type:
            sql += " AND j.visa_types LIKE %s"
            params.append(f"%{visa_type}%")
        if local_government_id:
            sql += " AND j.local_government_id = %s"
            params.append(local_government_id)
        sql += " ORDER BY j.created_at DESC"
        cursor.execute(sql, params)
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
        return {"items": items, "total": len(items)}
    finally:
        conn.close()

@router.get("/jobs/{job_id}")
async def get_job(job_id: int):
    """구인 상세"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            SELECT j.*, u.name as created_by_name, lg.name as lg_name
            FROM kwv_jobs j
            LEFT JOIN kwv_users u ON j.created_by = u.id
            LEFT JOIN kwv_local_governments lg ON j.local_government_id = lg.id
            WHERE j.id = %s
        """, (job_id,))
        item = cursor.fetchone()
        if not item:
            raise HTTPException(status_code=404, detail="구인 공고를 찾을 수 없습니다")
        for k, v in item.items():
            if isinstance(v, (datetime, date)):
                item[k] = v.isoformat()
        return item
    finally:
        conn.close()

@router.post("/admin/jobs")
async def create_job(request: Request, user: dict = Depends(get_current_user)):
    """구인 등록 (admin≥2)"""
    require_admin_level(user, 2)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            INSERT INTO kwv_jobs (title, description, local_government_id, visa_types,
                positions, salary, period, location, requirements, benefits,
                contact_name, contact_phone, contact_email, status, image_url, created_by)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            body['title'], body.get('description', ''),
            body.get('local_government_id') or None,
            body.get('visa_types', ''),
            body.get('positions', 1),
            body.get('salary', ''),
            body.get('period', ''),
            body.get('location', ''),
            body.get('requirements', ''),
            body.get('benefits', ''),
            body.get('contact_name', ''),
            body.get('contact_phone', ''),
            body.get('contact_email', ''),
            body.get('status', 'draft'),
            body.get('image_url', ''),
            int(user['sub'])
        ))
        conn.commit()
        return {"message": "구인 공고가 등록되었습니다", "id": cursor.lastrowid}
    finally:
        conn.close()

@router.put("/admin/jobs/{job_id}")
async def update_job(job_id: int, request: Request, user: dict = Depends(get_current_user)):
    """구인 수정 (admin≥2)"""
    require_admin_level(user, 2)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            UPDATE kwv_jobs SET title=%s, description=%s, local_government_id=%s,
                visa_types=%s, positions=%s, salary=%s, period=%s, location=%s,
                requirements=%s, benefits=%s, contact_name=%s, contact_phone=%s,
                contact_email=%s, status=%s, image_url=%s
            WHERE id=%s
        """, (
            body['title'], body.get('description', ''),
            body.get('local_government_id') or None,
            body.get('visa_types', ''),
            body.get('positions', 1),
            body.get('salary', ''),
            body.get('period', ''),
            body.get('location', ''),
            body.get('requirements', ''),
            body.get('benefits', ''),
            body.get('contact_name', ''),
            body.get('contact_phone', ''),
            body.get('contact_email', ''),
            body.get('status', 'draft'),
            body.get('image_url', ''),
            job_id
        ))
        conn.commit()
        return {"message": "구인 공고가 수정되었습니다"}
    finally:
        conn.close()

@router.delete("/admin/jobs/{job_id}")
async def delete_job(job_id: int, user: dict = Depends(get_current_user)):
    """구인 삭제 (admin≥3)"""
    require_admin_level(user, 3)
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("DELETE FROM kwv_jobs WHERE id = %s", (job_id,))
        conn.commit()
        return {"message": "구인 공고가 삭제되었습니다"}
    finally:
        conn.close()

@router.post("/jobs/{job_id}/apply")
async def apply_for_job(job_id: int, request: Request, user: dict = Depends(get_current_user)):
    """구인에 지원"""
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        # 구인 공고 확인
        cursor.execute("SELECT id, title, status FROM kwv_jobs WHERE id = %s", (job_id,))
        job = cursor.fetchone()
        if not job:
            raise HTTPException(status_code=404, detail="구인 공고를 찾을 수 없습니다")
        if job['status'] != 'active':
            raise HTTPException(status_code=400, detail="현재 지원 가능한 공고가 아닙니다")
        body = await request.json()
        user_id = int(user['sub'])
        # 알림 생성 (관리자에게)
        try:
            user_name = user.get('name', user.get('email', ''))
            cursor.execute("""
                INSERT INTO kwv_notifications (user_id, type, title, message)
                SELECT id, 'info', %s, %s FROM kwv_users WHERE user_type = 'admin' AND admin_level >= 2 AND is_active = TRUE
            """, (
                f"구직 지원: {job['title']}",
                f"{user_name}님이 '{job['title']}' 공고에 지원했습니다. 메시지: {body.get('message', '')}"
            ))
        except:
            pass
        conn.commit()
        return {"message": "지원이 완료되었습니다"}
    finally:
        conn.close()

# ==================== 관리자 관리 API ====================

@router.get("/admin/admins")
async def list_admins(user: dict = Depends(get_current_user)):
    """관리자 목록 (admin≥3)"""
    require_admin_level(user, 3)
    conn = get_kwv_db_connection()
    if not conn:
        return {"items": []}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            SELECT id, email, name, phone, organization, admin_level, is_active, is_approved,
                   user_type, created_at, last_login_at as last_login
            FROM kwv_users
            WHERE user_type = 'admin'
            ORDER BY admin_level DESC, created_at ASC
        """)
        items = cursor.fetchall()
        for item in items:
            for k, v in item.items():
                if isinstance(v, (datetime, date)):
                    item[k] = v.isoformat()
        return {"items": items, "total": len(items)}
    finally:
        conn.close()

@router.post("/admin/admins")
async def create_admin(request: Request, user: dict = Depends(get_current_user)):
    """관리자 추가 (admin≥9 = super admin only)"""
    require_admin_level(user, 9)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        # 이메일 중복 체크
        cursor.execute("SELECT id FROM kwv_users WHERE email = %s", (body['email'],))
        if cursor.fetchone():
            raise HTTPException(status_code=400, detail="이미 등록된 이메일입니다")
        hashed_pw = hash_password(body.get('password', DEFAULT_PASSWORD))
        cursor.execute("""
            INSERT INTO kwv_users (email, name, password, phone, organization,
                user_type, admin_level, is_active, is_approved)
            VALUES (%s, %s, %s, %s, %s, 'admin', %s, 1, 1)
        """, (
            body['email'], body['name'], hashed_pw,
            body.get('phone', ''),
            body.get('organization', ''),
            body.get('admin_level', 1)
        ))
        conn.commit()
        return {"message": "관리자가 추가되었습니다", "id": cursor.lastrowid}
    finally:
        conn.close()

@router.put("/admin/admins/{admin_id}")
async def update_admin(admin_id: int, request: Request, user: dict = Depends(get_current_user)):
    """관리자 정보/등급 수정 (admin≥9)"""
    require_admin_level(user, 9)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        fields = []
        params = []
        for field in ['name', 'phone', 'organization', 'admin_level', 'is_active', 'is_approved']:
            if field in body:
                fields.append(f"{field} = %s")
                params.append(body[field])
        # 승인 시 approved_at 설정
        if body.get('is_approved') or (body.get('admin_level') and int(body.get('admin_level', 0)) > 0):
            fields.append("is_approved = TRUE")
            fields.append("approved_at = NOW()")
        if 'password' in body and body['password']:
            fields.append("password = %s")
            params.append(hash_password(body['password']))
        if not fields:
            return {"message": "변경할 내용이 없습니다"}
        params.append(admin_id)
        cursor.execute(f"UPDATE kwv_users SET {', '.join(fields)} WHERE id = %s AND user_type = 'admin'", params)
        conn.commit()
        return {"message": "관리자 정보가 수정되었습니다"}
    finally:
        conn.close()

@router.delete("/admin/admins/{admin_id}")
async def delete_admin(admin_id: int, user: dict = Depends(get_current_user)):
    """관리자 삭제 (admin≥9, 자기 자신 삭제 불가)"""
    require_admin_level(user, 9)
    if int(user['sub']) == admin_id:
        raise HTTPException(status_code=400, detail="자기 자신은 삭제할 수 없습니다")
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("UPDATE kwv_users SET is_active = 0 WHERE id = %s AND user_type = 'admin'", (admin_id,))
        conn.commit()
        return {"message": "관리자가 삭제되었습니다"}
    finally:
        conn.close()

# ==================== 테마 설정 API ====================

@router.get("/admin/theme")
async def get_theme_settings(user: dict = Depends(get_current_user)):
    """테마 설정 조회"""
    conn = get_kwv_db_connection()
    if not conn:
        return {}
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("SELECT setting_key, setting_value FROM kwv_system_settings WHERE setting_key LIKE 'theme_%'")
        rows = cursor.fetchall()
        result = {}
        for row in rows:
            result[row['setting_key']] = row['setting_value']
        return result
    finally:
        conn.close()

@router.put("/admin/theme")
async def update_theme_settings(request: Request, user: dict = Depends(get_current_user)):
    """테마 설정 저장 (admin≥3)"""
    require_admin_level(user, 3)
    body = await request.json()
    conn = get_kwv_db_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="DB 연결 실패")
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        for key in ['theme_preset', 'theme_header_bg_start', 'theme_header_bg_end', 'theme_menu_active_color']:
            if key in body:
                cursor.execute("""
                    INSERT INTO kwv_system_settings (setting_key, setting_value, setting_type, description, updated_by)
                    VALUES (%s, %s, 'string', %s, %s)
                    ON DUPLICATE KEY UPDATE setting_value = VALUES(setting_value), updated_by = VALUES(updated_by)
                """, (key, body[key], f'테마: {key}', int(user['sub'])))
        conn.commit()
        return {"message": "테마 설정이 저장되었습니다"}
    finally:
        conn.close()

# ==================== Health Check ====================

@router.get("/health")
async def health_check():
    """서버 상태 확인"""
    return {
        "status": "ok",
        "service": "KoreaWorkingVisa API",
        "version": "2.1.20260214",
        "mock_mode": MOCK_MODE,
        "jwt_available": JWT_AVAILABLE,
        "bcrypt_available": BCRYPT_AVAILABLE,
        "google_oauth_configured": bool(GOOGLE_CLIENT_ID)
    }

# ==================== RAG 문서 관리 ====================

@router.get("/admin/rag/documents")
async def get_rag_documents(
    category: str = None, search: str = None, active_only: bool = True,
    page: int = 1, limit: int = 20,
    user: dict = Depends(get_current_user)
):
    """RAG 문서 목록 조회"""
    conn = get_kwv_db_connection()
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        where = []
        params = []
        if active_only:
            where.append("d.is_active = 1")
        if category:
            where.append("d.category = %s")
            params.append(category)
        if search:
            where.append("(d.title LIKE %s OR d.content LIKE %s OR d.tags LIKE %s)")
            s = f"%{search}%"
            params.extend([s, s, s])
        where_sql = " AND ".join(where) if where else "1=1"
        offset = (page - 1) * limit

        cursor.execute(f"SELECT COUNT(*) as cnt FROM kwv_rag_documents d WHERE {where_sql}", params)
        total = cursor.fetchone()['cnt']

        cursor.execute(f"""
            SELECT d.*, u.name as creator_name
            FROM kwv_rag_documents d
            LEFT JOIN kwv_users u ON d.created_by = u.id
            WHERE {where_sql}
            ORDER BY d.updated_at DESC
            LIMIT %s OFFSET %s
        """, params + [limit, offset])
        docs = cursor.fetchall()
        for doc in docs:
            for k, v in doc.items():
                if hasattr(v, 'isoformat'):
                    doc[k] = v.isoformat()
        return {"success": True, "items": docs, "total": total, "page": page, "pages": (total + limit - 1) // limit}
    finally:
        conn.close()

@router.get("/admin/rag/documents/{doc_id}")
async def get_rag_document(doc_id: int, user: dict = Depends(get_current_user)):
    """RAG 문서 상세 조회"""
    conn = get_kwv_db_connection()
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            SELECT d.*, u.name as creator_name
            FROM kwv_rag_documents d
            LEFT JOIN kwv_users u ON d.created_by = u.id
            WHERE d.id = %s
        """, (doc_id,))
        doc = cursor.fetchone()
        if not doc:
            raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다")
        for k, v in doc.items():
            if hasattr(v, 'isoformat'):
                doc[k] = v.isoformat()
        return {"success": True, "data": doc}
    finally:
        conn.close()

@router.post("/admin/rag/documents")
async def create_rag_document(request: Request, user: dict = Depends(get_current_user)):
    """RAG 문서 등록"""
    admin_level = user.get('admin_level', 0)
    if admin_level < 2:
        raise HTTPException(status_code=403, detail="관리자 권한이 필요합니다")
    body = await request.json()
    title = body.get('title', '').strip()
    content = body.get('content', '').strip()
    if not title or not content:
        raise HTTPException(status_code=400, detail="제목과 내용은 필수입니다")
    conn = get_kwv_db_connection()
    try:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO kwv_rag_documents (title, category, content, content_summary, tags, language, created_by, updated_by)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            title, body.get('category', 'general'), content,
            body.get('content_summary', content[:200]),
            body.get('tags', ''), body.get('language', 'ko'),
            int(user['sub']), int(user['sub'])
        ))
        conn.commit()
        return {"success": True, "message": "문서가 등록되었습니다", "id": cursor.lastrowid}
    finally:
        conn.close()

@router.put("/admin/rag/documents/{doc_id}")
async def update_rag_document(doc_id: int, request: Request, user: dict = Depends(get_current_user)):
    """RAG 문서 수정"""
    admin_level = user.get('admin_level', 0)
    if admin_level < 2:
        raise HTTPException(status_code=403, detail="관리자 권한이 필요합니다")
    body = await request.json()
    conn = get_kwv_db_connection()
    try:
        cursor = conn.cursor()
        cursor.execute("""
            UPDATE kwv_rag_documents
            SET title=%s, category=%s, content=%s, content_summary=%s, tags=%s, language=%s, is_active=%s, updated_by=%s
            WHERE id=%s
        """, (
            body.get('title'), body.get('category', 'general'), body.get('content'),
            body.get('content_summary', body.get('content', '')[:200]),
            body.get('tags', ''), body.get('language', 'ko'),
            body.get('is_active', 1), int(user['sub']), doc_id
        ))
        conn.commit()
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다")
        return {"success": True, "message": "문서가 수정되었습니다"}
    finally:
        conn.close()

@router.post("/admin/rag/documents/upload")
async def upload_rag_document(
    file: UploadFile = File(...),
    category: str = Form("general"),
    language: str = Form("ko"),
    user: dict = Depends(get_current_user)
):
    """RAG 문서 파일 업로드 (PDF, DOCX, XLSX, PPTX)"""
    admin_level = user.get('admin_level', 0)
    if admin_level < 2:
        raise HTTPException(status_code=403, detail="관리자 권한이 필요합니다")

    # MIME 타입 및 확장자 검증
    allowed_types = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    }
    allowed_exts = {'pdf', 'docx', 'xlsx', 'pptx'}
    file_ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in (file.filename or '') else ''
    ext = allowed_types.get(file.content_type) or (file_ext if file_ext in allowed_exts else None)
    if not ext:
        raise HTTPException(status_code=400, detail="허용되지 않는 파일 형식입니다. (PDF, DOCX, XLSX, PPTX만 가능)")

    file_data = await file.read()
    file_size = len(file_data)
    if file_size > 20 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="파일 크기가 20MB를 초과합니다")
    extracted_text = ""
    import tempfile
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name

        if ext == 'pdf':
            import PyPDF2
            reader = PyPDF2.PdfReader(tmp_path)
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text.strip())
            extracted_text = "\n\n".join(pages)

        elif ext == 'docx':
            import docx
            doc = docx.Document(tmp_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            extracted_text = "\n".join(paragraphs)

        elif ext == 'xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(tmp_path, read_only=True, data_only=True)
            sheets_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    sheets_text.append(f"[{sheet_name}]\n" + "\n".join(rows))
            wb.close()
            extracted_text = "\n\n".join(sheets_text)

        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(tmp_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    slides_text.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
            extracted_text = "\n\n".join(slides_text)

    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.remove(tmp_path)

    if not extracted_text.strip():
        raise HTTPException(status_code=400, detail="파일에서 텍스트를 추출할 수 없습니다")

    # 원본 파일 저장
    rag_dir = os.path.join(LOCAL_UPLOAD_DIR, "rag")
    os.makedirs(rag_dir, exist_ok=True)
    save_filename = f"{int(datetime.utcnow().timestamp())}_{uuid.uuid4().hex[:8]}.{ext}"
    save_path = os.path.join(rag_dir, save_filename)
    with open(save_path, 'wb') as f:
        f.write(file_data)
    file_url = f"/api/kwv/uploads/rag/{save_filename}"

    # DB 저장
    title = os.path.splitext(file.filename)[0]
    content_summary = extracted_text[:200]
    conn = get_kwv_db_connection()
    try:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO kwv_rag_documents
            (title, category, content, content_summary, tags, language, created_by, updated_by, file_path, file_name, file_size)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            title, category, extracted_text, content_summary,
            ext, language, int(user['sub']), int(user['sub']),
            file_url, file.filename, file_size
        ))
        conn.commit()
        return {
            "success": True,
            "message": f"파일에서 텍스트를 추출하여 문서로 등록했습니다 ({len(extracted_text)}자)",
            "id": cursor.lastrowid,
            "extracted_length": len(extracted_text)
        }
    finally:
        conn.close()

@router.delete("/admin/rag/documents/{doc_id}")
async def delete_rag_document(doc_id: int, user: dict = Depends(get_current_user)):
    """RAG 문서 삭제"""
    admin_level = user.get('admin_level', 0)
    if admin_level < 2:
        raise HTTPException(status_code=403, detail="관리자 권한이 필요합니다")
    conn = get_kwv_db_connection()
    try:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM kwv_rag_documents WHERE id=%s", (doc_id,))
        conn.commit()
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다")
        return {"success": True, "message": "문서가 삭제되었습니다"}
    finally:
        conn.close()

@router.get("/rag/search")
async def search_rag_documents(q: str, limit: int = 5):
    """RAG 문서 검색 (챗봇용, 인증 불필요)"""
    conn = get_kwv_db_connection()
    try:
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        cursor.execute("""
            SELECT id, title, category, content, content_summary, tags
            FROM kwv_rag_documents
            WHERE is_active = 1 AND (title LIKE %s OR content LIKE %s OR tags LIKE %s)
            ORDER BY updated_at DESC
            LIMIT %s
        """, (f"%{q}%", f"%{q}%", f"%{q}%", limit))
        docs = cursor.fetchall()
        return {"success": True, "items": docs}
    finally:
        conn.close()
